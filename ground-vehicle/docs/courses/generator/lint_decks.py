"""
lint_decks.py - static checks on the generated decks.

Without a renderer installed we cannot look at the slides, so this estimates
instead: it works out roughly how tall each block of text will be once it
wraps, and reports anything that will not fit in the shape it was put in, or
that runs off the edge of the slide.

The estimate is deliberately pessimistic - it would rather warn about a slide
that is fine than stay quiet about one that is not.

    python lint_decks.py
"""

import glob
import math
import os
import sys

from pptx import Presentation
from pptx.util import Emu, Pt

EMU_PER_PT = 12700
SLIDE_W_PT = 13.333 * 72
SLIDE_H_PT = 7.5 * 72

# Calibri and Consolas average advance width, as a fraction of point size.
CHAR_W = {"Consolas": 0.55, "Calibri": 0.48, None: 0.48}

# Tolerance before we complain, in points.
SLOP_PT = 6.0


def estimate_height_pt(frame, width_pt):
    """Roughly how tall this text frame wants to be, in points."""
    total = 0.0
    for para in frame.paragraphs:
        text = "".join(run.text for run in para.runs) or para.text or ""
        size = None
        font = None
        for run in para.runs:
            if run.font.size:
                size = run.font.size.pt
            if run.font.name:
                font = run.font.name
            break
        # A paragraph with no runs (a blank spacer line) still occupies a line,
        # at whatever its paragraph-level size says - or the theme default.
        if size is None and para.font.size is not None:
            size = para.font.size.pt
        if font is None and para.font.name is not None:
            font = para.font.name
        if size is None:
            size = 18.0

        indent = para.level * 18.0
        usable = max(width_pt - indent, 24.0)
        per_line = max(int(usable / (CHAR_W.get(font, 0.48) * size)), 1)

        lines = max(math.ceil(len(text) / per_line), 1) if text else 1
        spacing = para.line_spacing if isinstance(para.line_spacing, float) else 1.0
        after = para.space_after.pt if para.space_after is not None else 0.0

        total += lines * size * 1.22 * spacing + after
    return total


def overlapping_text(shapes, slide_h):
    """
    Pairs of text-bearing shapes that sit on top of each other.

    The height estimate above catches text that will not fit its own box. It
    says nothing about two boxes that each fit and are drawn in the same
    place - a note panel over a code panel, a heading over a row of labels -
    which is the other way a slide goes wrong. Only substantial overlaps are
    reported, so a caption tucked under a picture edge stays quiet.
    """
    boxes = []
    for shape in shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        try:
            left, top = shape.left, shape.top
            width, height = shape.width, shape.height
        except (AttributeError, TypeError):
            continue
        if None in (left, top, width, height) or width <= 0 or height <= 0:
            continue
        # A footer strip runs the width of the slide by design.
        if top > slide_h - Emu(500000):
            continue
        boxes.append((left, top, width, height, text))

    found = []
    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            ax, ay, aw, ah, atext = boxes[i]
            bx, by, bw, bh, btext = boxes[j]
            over_w = min(ax + aw, bx + bw) - max(ax, bx)
            over_h = min(ay + ah, by + bh) - max(ay, by)
            if over_w <= 0 or over_h <= 0:
                continue
            area = over_w * over_h
            smaller = min(aw * ah, bw * bh)
            if smaller and area / smaller > 0.45:
                found.append("{!r} over {!r}".format(atext[:26], btext[:26]))
    return found


def check_deck(path):
    prs = Presentation(path)
    issues = []
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for index, slide in enumerate(prs.slides, 1):
        for clash in overlapping_text(slide.shapes, slide_h):
            issues.append((index, "text shapes overlap", clash))

        for shape in slide.shapes:
            name = shape.shape_type

            # Off the edge of the slide?
            try:
                left, top = shape.left, shape.top
                width, height = shape.width, shape.height
            except (AttributeError, TypeError):
                continue
            if left is None or top is None or width is None or height is None:
                continue

            if left < -Emu(1000) or top < -Emu(1000):
                issues.append((index, "off slide (negative)", str(name)))
            if left + width > slide_w + Emu(20000):
                issues.append((index, "runs off the right edge",
                               "{:.2f} in".format((left + width) / 914400)))
            if top + height > slide_h + Emu(20000):
                issues.append((index, "runs off the bottom",
                               "{:.2f} in".format((top + height) / 914400)))

            if not shape.has_text_frame:
                continue
            frame = shape.text_frame
            if not frame.text.strip():
                continue
            # The title-slide credit block is chrome, sized like the
            # footer, and is allowed to sit below the body floor.
            if "Porpoise Robotics" in frame.text and "President" in frame.text:
                continue

            inset_l = frame.margin_left.pt if frame.margin_left is not None else 7.2
            inset_r = frame.margin_right.pt if frame.margin_right is not None else 7.2
            inset_t = frame.margin_top.pt if frame.margin_top is not None else 3.6
            inset_b = frame.margin_bottom.pt if frame.margin_bottom is not None else 3.6

            width_pt = width / EMU_PER_PT - inset_l - inset_r
            height_pt = height / EMU_PER_PT - inset_t - inset_b
            if width_pt <= 0 or height_pt <= 0:
                continue

            wanted = estimate_height_pt(frame, width_pt)
            if wanted > height_pt + SLOP_PT:
                preview = frame.text.strip().replace("\n", " / ")[:58]
                issues.append((index,
                               "text overflows by {:.0f} pt".format(wanted - height_pt),
                               preview))

    return issues


def main():
    root = os.path.normpath(os.path.join(os.path.dirname(os.path.abspath(__file__)), ".."))
    paths = sorted(glob.glob(os.path.join(root, "*", "*.pptx")))

    total = 0
    for path in paths:
        issues = check_deck(path)
        label = os.path.join(os.path.basename(os.path.dirname(path)),
                             os.path.basename(path))
        if issues:
            print("\n{}  ({} issues)".format(label, len(issues)))
            for slide_no, kind, detail in issues:
                print("   slide {:>3}  {:<28} {}".format(slide_no, kind, detail))
            total += len(issues)
        else:
            print("{}  clean".format(label))

    print("\n{} issues across {} decks".format(total, len(paths)))
    return 1 if total else 0


if __name__ == "__main__":
    sys.exit(main())
