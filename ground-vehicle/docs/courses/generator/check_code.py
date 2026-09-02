"""
check_code.py - stop the slides drifting away from the code.

Every code listing on a slide is supposed to be a real excerpt from a real file
in this repository. It is very easy for that to stop being true: somebody
improves a function, and a slide somewhere quietly goes on teaching last
month's version. That is exactly what happened to the P2 and P3 decks.

This walks every built deck, pulls out every line set in the code font, and
checks it appears verbatim in one of the sources under ground-vehicle/src/.
Anything it cannot find is either drift, or an illustrative line that belongs
in ALLOWED below.

    python check_code.py            report only
    python check_code.py -v         also list which file each line matched

Exit status is 1 if anything is unaccounted for, so it can gate a commit.
"""

import glob
import os
import re
import sys

from pptx import Presentation

HERE = os.path.dirname(os.path.abspath(__file__))
COURSES = os.path.normpath(os.path.join(HERE, ".."))
SRC = os.path.normpath(os.path.join(COURSES, "..", "..", "src"))

CODE_FONT = "Consolas"
PANEL_NAME = "CODE_PANEL"

# Lines that are deliberately NOT quotes from the repository: generic C taught
# as language, the millis() pattern written in the abstract, and code that is
# deliberately wrong so students can see the bug.
ALLOWED = {
    # generic C and Arduino, taught as language rather than as our code
    "for (int i = 0; i < 32; i++) {",
    "strip.setPixelColor(i, colour);",
    "strip.setPixelColor(0, white);",
    "strip.setPixelColor(1, white);",
    "strip.setPixelColor(2, white);",
    "while (!Ps3.isConnected()) {",
    "showWaitingLights();",
    "do {",
    "reading = analogRead(PIN);",
    "} while (reading > LIMIT);",
    "map(value, fromLow, fromHigh, toLow, toHigh)",
    "size = map(abs(v), 20, 127, MOTOR_MIN, maxSpeed);",
    "size = map(abs(v), 60, 511, MOTOR_MIN, maxSpeed);",
    "if (abs(v) < DEADZONE)",
    "return (v > 0) ? size",
    ": -size;",
    "int size = map(abs(v),",
    "DEADZONE, 127,",
    "DEADZONE, 511,",
    "MOTOR_MIN, maxSpeed);",

    # the millis() pattern, written generically
    "unsigned long lastBlink = 0;",
    "if (millis() - lastBlink >= INTERVAL) {",
    "if (now - last_update < INTERVAL) return;",
    "last_update = now;",
    "if (millis() - lastTime >= INTERVAL) { lastTime = millis(); ...do the thing... }",

    # deliberately wrong, shown so students can see the bug
    "if (Ps3.data.button.square) {",
    "if (myController->x()) {",
    "lightsOn = !lightsOn;",
    "current += (target - current) * RAMP_FACTOR;",

    # edge detection, written as the idea rather than as the function
    "bool isNewPress = isDown && !wasDown;",
    "wasDown = isDown;",

    # allowlist calls, shown without their surrounding function
    "uni_bt_allowlist_remove_all();",
    "uni_bt_allowlist_add_addr(myControllerAddress);",
    "uni_bt_allowlist_add_addr(address);",
    "uni_bt_allowlist_set_enabled(true);",
    "BP32.setup(&onConnectedController, &onDisconnectedController);",

    # worked arithmetic showing what map() returns, not lines of the program
    "map(0, 0, 127, 0, 255) -> 0",
    "map(63, 0, 63, 0, 255) -> 127",
    "map(127, 0, 127, 0, 255) -> 255",
    "map(0, 0, 511, 0, 255) -> 0",
    "map(255, 0, 255, 0, 255) -> 127",
    "map(511, 0, 511, 0, 255) -> 255",

    # formulae and register names set in the code font
    "duty = microseconds * 65536 / 20000",
    "V = I x R", "I = V / R", "R = V / I",
    "P = I x V", "I = P / V", "V = P / I",
}

COMMENT_ONLY = re.compile(r"^\s*(//|/\*|\*)")
TRAILING_COMMENT = re.compile(r"\s*//.*$")


def normalise(line):
    return re.sub(r"\s+", " ", line).strip()


def load_sources():
    corpus = {}
    for pattern in ("**/*.ino", "**/*.h"):
        for path in glob.glob(os.path.join(SRC, pattern), recursive=True):
            rel = os.path.relpath(path, SRC).replace("\\", "/")
            with open(path, encoding="utf-8", errors="replace") as handle:
                corpus[rel] = {normalise(l) for l in handle if l.strip()}
    return corpus


def code_lines(deck_path):
    """
    Every line inside a tagged code panel, with its slide number and its
    position in that panel, so a line that was wrapped to fit the slide can be
    rejoined with the one after it before matching.
    """
    prs = Presentation(deck_path)
    out = []
    for number, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame or shape.name != PANEL_NAME:
                continue
            panel = [normalise("".join(r.text for r in p.runs))
                     for p in shape.text_frame.paragraphs]
            for index, text in enumerate(panel):
                if text:
                    out.append((number, text, panel, index))
    return out


def matches(line, corpus):
    """Is this line, verbatim, in some source file? Returns the file or None."""
    for name, lines in corpus.items():
        for source in lines:
            if source == line:
                return name
            # The source line may carry a trailing comment the slide dropped.
            if normalise(TRAILING_COMMENT.sub("", source)) == line:
                return name
    return None


def matches_wrapped(panel, index, corpus, span=3):
    """
    A long source line is sometimes broken over two or three slide lines so it
    fits the panel. Join this line with the ones after it and try again.
    """
    for extra in range(1, span + 1):
        if index + extra >= len(panel):
            break
        joined = normalise(" ".join(panel[index:index + extra + 1]))
        joined = normalise(TRAILING_COMMENT.sub("", joined))
        hit = matches(joined, corpus)
        if hit:
            return hit, extra
    return None, 0


def main():
    verbose = "-v" in sys.argv
    corpus = load_sources()
    if not corpus:
        sys.exit("No sources found under %s" % SRC)

    print("checked against %d source files under ground-vehicle/src/\n" % len(corpus))

    total_lines = 0
    unmatched = []

    for deck in sorted(glob.glob(os.path.join(COURSES, "*", "*.pptx"))):
        label = os.path.join(os.path.basename(os.path.dirname(deck)),
                             os.path.basename(deck))
        skip_until = -1
        for number, raw, panel, index in code_lines(deck):
            if index <= skip_until:
                continue

            # A comment on its own is the slide talking, not the program.
            if COMMENT_ONLY.match(raw):
                continue

            # A trailing comment is often shortened to fit the slide width.
            # The CODE has to match; the commentary does not.
            line = normalise(TRAILING_COMMENT.sub("", raw))

            if not line or line in {"...", "{", "}", "};", "});"}:
                continue
            total_lines += 1

            if line in ALLOWED:
                continue

            hit = matches(line, corpus)
            if not hit:
                hit, consumed = matches_wrapped(panel, index, corpus)
                if hit:
                    skip_until = index + consumed
            if hit:
                if verbose:
                    print("  ok   %-44s s%-3d %-44s  %s"
                          % (label, number, line[:44], hit))
                continue

            unmatched.append((label, number, line))

    print("%d code lines checked" % total_lines)

    if unmatched:
        print("\n%d line(s) not found in any source file:\n" % len(unmatched))
        for label, number, line in unmatched:
            print("  %-46s slide %-3d %s" % (label, number, line))
        print("\nEither the slide has drifted from the code, or the line is")
        print("illustrative and belongs in ALLOWED in this file.")
        return 1

    print("no drift: every code line on every slide is in the repository")
    return 0


if __name__ == "__main__":
    sys.exit(main())
