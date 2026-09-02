"""
diagrams.py - the drawn figures used in the Pathfinder course decks.

Everything here is built from native PowerPoint autoshapes and text boxes, so
each figure stays editable, scales without going fuzzy, and prints sharp in
black and white. Nothing is a bitmap.

Each function takes a Deck, adds one slide, and returns it.
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from slidelib import (AMBER, BODY_FONT, CODE_BG, CODE_FONT, CONTENT_W,
                      EMU_PER_PT, GREY, INK, MARGIN_L, MIN_SMALL_PT, NAVY,
                      MIN_CODE_PT, PANEL, RED, RULE, TEAL, WHITE, BODY_TOP,
                      measure_pt,
                      _set_fitted, _set_text)

LIGHT_TEAL = RGBColor(0xD6, 0xEC, 0xEE)
LIGHT_AMBER = RGBColor(0xFA, 0xE6, 0xC8)
LIGHT_RED = RGBColor(0xF6, 0xD8, 0xD8)
LIGHT_GREY = RGBColor(0xE8, 0xEC, 0xF0)


# ===================================================================
# small drawing helpers
# ===================================================================

def _box(slide, left, top, width, height, text, *, fill=WHITE, edge=NAVY,
         size=13, bold=False, color=INK, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         font=BODY_FONT, edge_w=1.25, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE):
    box = slide.shapes.add_shape(shape, left, top, width, height)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        box.adjustments[0] = 0.12
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = edge
    box.line.width = Pt(edge_w)
    box.shadow.inherit = False

    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = anchor
    # A box holding CODE follows the code floor, not the prose floor - the
    # same rule the full-slide listings use.
    floor = MIN_CODE_PT if font == CODE_FONT else MIN_SMALL_PT
    _set_fitted(frame, text if isinstance(text, list) else [text],
                width=width, height=height, size=size, floor=floor,
                bold=bold, color=color, align=align, space_after=0,
                font=font)
    return box


def _label(slide, left, top, width, text, *, size=12, bold=False, color=INK,
           align=PP_ALIGN.LEFT, font=BODY_FONT, height=None):
    # A label with no explicit height gets one that suits how many lines it
    # holds, so a multi-line note is never squeezed into a single-line box.
    # Size it for the font it will ACTUALLY render at: the house minimum wins
    # over whatever the caller asked for, so a 12pt request still gets a box
    # tall enough for 18pt text.
    effective = max(size, MIN_SMALL_PT)
    if height is None:
        # Measure the text properly rather than counting list items: a short
        # label in a narrow box still wraps onto a second line at 18pt, and
        # counting items would size the box for one.
        items = text if isinstance(text, list) else [text]
        wanted = measure_pt(items, width / EMU_PER_PT, effective,
                            font=font, space_after=1)
        height = Inches(max(0.34, (wanted + 6) / 72.0))
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_top = 0
    frame.margin_right = 0
    frame.margin_bottom = 0
    _set_fitted(frame, text if isinstance(text, list) else [text],
                width=width, height=height, size=size, floor=MIN_SMALL_PT,
                bold=bold, color=color, align=align, space_after=1, font=font)
    return box


def _arrow(slide, x1, y1, x2, y2, *, color=NAVY, width=1.75, dashed=False):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if dashed:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    # python-pptx has no arrowhead API, so set it on the XML directly.
    tail = line.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    head = tail.makeelement(qn("a:headEnd"), {})
    end = tail.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med",
                                             "len": "med"})
    tail.append(head)
    tail.append(end)
    return line


def _plain_line(slide, x1, y1, x2, y2, *, color=GREY, width=1.0):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


# ===================================================================
# 1. the three parts of a program
# ===================================================================

def program_structure(deck, title="Every program has the same three parts", speaker=None):
    slide = deck.blank(title, speaker=speaker or [
        "Hold this shape up for the whole course. Every program they ever open has these three parts, and knowing which part a line lives in tells you WHEN it runs.",
        "Ask before you explain: which of these three runs more than once? Let somebody get it wrong first - the answer sticks better.",
        "The first part has no keyword in C++. It is just the top of the file. Naming it anyway gives them somewhere to file #includes and constants.",
        "Common misunderstanding: students think setup() runs every time round. Say plainly that it runs once, at power-up, and never again until reset.",
    ])

    # Named the way the course names them: INITIALIZATION, SETUP, LOOP. The
    # first has no keyword of its own in C++ - it is simply the top of the
    # file - but giving it a name makes the three-part shape easy to hold on to.
    parts = [
        ("1.  INITIALIZATION", "Libraries, constants, variables.\nThe top of the file.",
         "#include <Adafruit_NeoPixel.h>\nconst int LED_PIN = 5;", LIGHT_GREY),
        ("2.  SETUP", "setup() runs ONCE, at power-up.\nGet the hardware ready.",
         "void setup() {\n  pinMode(LED_PIN, OUTPUT);\n}", LIGHT_TEAL),
        ("3.  LOOP", "loop() runs OVER AND OVER,\nuntil the power goes off.",
         "void loop() {\n  digitalWrite(LED_PIN, HIGH);\n  delay(1000);\n}", LIGHT_AMBER),
    ]

    col_w = Inches(3.85)
    gap = Inches(0.32)
    top = BODY_TOP + Inches(0.15)

    for index, (heading, blurb, code, fill) in enumerate(parts):
        left = MARGIN_L + (col_w + gap) * index

        _box(slide, left, top, col_w, Inches(0.5), heading, fill=fill,
             edge=NAVY, size=14, bold=True, color=NAVY)

        _label(slide, left + Inches(0.1), top + Inches(0.68), col_w - Inches(0.2),
               blurb.split("\n"), size=13, color=INK)

        _box(slide, left, top + Inches(1.55), col_w, Inches(1.5),
             code.split("\n"), fill=CODE_BG, edge=RULE, size=11,
             font=CODE_FONT, shape=MSO_SHAPE.RECTANGLE, edge_w=0.75,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

        if index < 2:
            _arrow(slide, left + col_w + Inches(0.05), top + Inches(1.6),
                   left + col_w + gap - Inches(0.05), top + Inches(1.6))

    # The loop arrow, curling back on itself.
    loop_left = MARGIN_L + (col_w + gap) * 2
    _label(slide, loop_left, top + Inches(3.2), col_w,
           "and round again, thousands of times a second", size=12,
           color=AMBER, align=PP_ALIGN.CENTER, bold=True)

    deck._note(slide,
               "Ingredients, then preparation, then cooking. A 500-line vehicle "
               "program has exactly this shape - there is just more in each part.",
               "info")
    return slide


# ===================================================================
# 2. system block diagram
# ===================================================================

def system_block(deck, title="How a command gets from your thumb to a wheel",
                 controller="PS3 controller", link="Bluetooth", speaker=None):
    slide = deck.blank(title, speaker=speaker or [
        "Walk it left to right with a real vehicle in your hand, touching each box as you name it.",
        "The point of the diagram is that nothing here is magic - every arrow is a wire or a radio link you can point at.",
        "Ask where THEIR program sits. The answer is the ESP32 box, and everything either side of it is fixed hardware they inherit.",
        "Come back to this slide whenever somebody is stuck: naming which box the fault is in is most of the debugging.",
    ])

    top = BODY_TOP + Inches(0.5)
    box_h = Inches(0.95)

    stages = [
        (controller, "you move a stick", LIGHT_TEAL, Inches(2.1)),
        (link, "radio, 2.4 GHz", WHITE, Inches(1.6)),
        ("ESP32", "your program runs here", LIGHT_TEAL, Inches(1.9)),
    ]

    left = MARGIN_L
    centres = []
    for name, blurb, fill, width in stages:
        _box(slide, left, top, width, box_h, name, fill=fill, size=15, bold=True,
             color=NAVY)
        _label(slide, left, top + box_h + Inches(0.06), width, blurb, size=11,
               color=GREY, align=PP_ALIGN.CENTER)
        centres.append((left, width))
        left += width + Inches(0.45)

    for index in range(len(stages) - 1):
        box_left, box_w = centres[index]
        _arrow(slide, box_left + box_w, top + Emu(int(box_h / 2)),
               box_left + box_w + Inches(0.45), top + Emu(int(box_h / 2)))

    # The ESP32 fans out to three kinds of output.
    esp_left, esp_w = centres[2]
    fan_left = esp_left + esp_w + Inches(0.45)
    outputs = [
        ("DRV8871 drivers", "4 motors", "12 13 16 17 18 19 22 23", LIGHT_AMBER),
        ("WS2812B strip", "32 LEDs", "GPIO 5", LIGHT_AMBER),
        ("Servo headers", "up to 4 servos", "25 26 27 14", LIGHT_AMBER),
    ]

    out_w = Inches(3.3)
    out_h = Inches(0.82)
    out_top = BODY_TOP + Inches(0.1)

    for index, (name, what, pins) in enumerate([(a, b, c) for a, b, c, _ in outputs]):
        y = out_top + (out_h + Inches(0.42)) * index
        _box(slide, fan_left + Inches(0.45), y, out_w, out_h,
             [name, what], fill=LIGHT_AMBER, edge=AMBER, size=13, bold=True,
             color=NAVY)
        _label(slide, fan_left + Inches(0.45), y + out_h + Inches(0.03), out_w,
               pins, size=10, color=GREY, align=PP_ALIGN.CENTER, font=CODE_FONT)
        _arrow(slide, fan_left, top + Emu(int(box_h / 2)),
               fan_left + Inches(0.45), y + Emu(int(out_h / 2)), color=AMBER)

    deck._note(slide,
               "Everything the vehicle can do happens on one of those three "
               "outputs. The rest of the program is deciding what to send.",
               "info")
    return slide


# ===================================================================
# 3. PWM duty cycle
# ===================================================================

def pwm_duty(deck, title="Pulse width modulation: how a pin makes a half speed", speaker=None):
    slide = deck.blank(title, speaker=speaker or [
        "Draw the square wave on the board first, then reveal the slide. It lands better if they have watched it being built.",
        "The idea to land: the pin is only ever fully on or fully off. Half speed is half the TIME on, not half the voltage.",
        "Analogy that works: a light switch flicked very fast. Too fast to see flicker, so the room looks half lit.",
        "Duty cycle is a percentage, which is why the maths lesson and the motor lesson are the same lesson.",
    ])

    left = MARGIN_L + Inches(1.55)
    width = Inches(8.2)
    row_h = Inches(0.98)
    top = BODY_TOP + Inches(0.10)
    cycles = 4

    rows = [
        ("25% duty", 0.25, "quarter power", "ledcWrite(pin, 64)"),
        ("50% duty", 0.50, "half power", "ledcWrite(pin, 128)"),
        ("100% duty", 1.00, "full power", "ledcWrite(pin, 255)"),
    ]

    for index, (name, fraction, meaning, call) in enumerate(rows):
        y = top + (row_h + Inches(0.32)) * index
        high = y
        low = y + Inches(0.62)

        _label(slide, MARGIN_L, y + Inches(0.12), Inches(1.45), name, size=14,
               bold=True, color=NAVY)

        # Baseline
        _plain_line(slide, left, low, left + width, low, color=RULE, width=1.0)

        cycle_w = Emu(int(width / cycles))
        on_w = Emu(int(cycle_w * fraction))

        for c in range(cycles):
            x0 = left + cycle_w * c
            if on_w > 0:
                _plain_line(slide, x0, high, x0 + on_w, high, color=TEAL, width=2.5)
                _plain_line(slide, x0, low, x0, high, color=TEAL, width=2.5)
                if fraction < 1.0:
                    _plain_line(slide, x0 + on_w, high, x0 + on_w, low,
                                color=TEAL, width=2.5)
            if fraction < 1.0:
                _plain_line(slide, x0 + on_w, low, x0 + cycle_w, low,
                            color=TEAL, width=2.5)

        _label(slide, left + width + Inches(0.12), y + Inches(0.02),
               Inches(2.6), [meaning, call], size=11, color=GREY)

    _label(slide, left, top + Inches(4.02), width,
           "one cycle = 1/20000 s; this picture repeats 5000 times a second",
           size=11, color=GREY, align=PP_ALIGN.CENTER)

    deck._note(slide,
               "The pin is only ever fully on or fully off. Speed comes from "
               "the FRACTION of each cycle it spends on - the duty cycle.",
               "info")
    return slide


# ===================================================================
# 4. H-bridge
# ===================================================================

def h_bridge(deck, title="The H-bridge: four switches, four things a motor can do", speaker=None):
    slide = deck.blank(title, speaker=speaker or [
        "Four switches, and only the two diagonal pairs are useful. Trace each path with a finger on the projected slide.",
        "Ask what happens if you close both switches on one side. Somebody will work out that it shorts the supply - which is exactly why the driver chip will not let you.",
        "Forward and reverse are the same circuit with the current running the other way. That is the whole trick, and it is why direction costs nothing extra in the program.",
        "Brake and coast are worth ten seconds each here; they come back properly in the advanced course.",
    ])

    states = [
        ("A high, B low", "FORWARD", "driven one way", LIGHT_TEAL, TEAL),
        ("A low, B high", "REVERSE", "driven the other way", LIGHT_TEAL, TEAL),
        ("both low", "COAST", "terminals floating,\nmotor freewheels", LIGHT_GREY, GREY),
        ("both high", "BRAKE", "terminals shorted,\nmotor stops hard", LIGHT_AMBER, AMBER),
    ]

    col_w = Inches(2.85)
    gap = Inches(0.28)
    top = BODY_TOP + Inches(0.35)

    for index, (pins, name, meaning, fill, edge) in enumerate(states):
        left = MARGIN_L + (col_w + gap) * index

        _box(slide, left, top, col_w, Inches(0.48), pins, fill=WHITE, edge=RULE,
             size=13, font=CODE_FONT, shape=MSO_SHAPE.RECTANGLE, edge_w=0.75)

        _box(slide, left, top + Inches(0.55), col_w, Inches(0.72), name,
             fill=fill, edge=edge, size=18, bold=True, color=NAVY)

        _label(slide, left, top + Inches(1.4), col_w, meaning.split("\n"),
               size=13, color=INK, align=PP_ALIGN.CENTER)

    # The bridge itself, drawn once underneath.
    bx = MARGIN_L + Inches(3.4)
    by = top + Inches(2.35)
    bw = Inches(6.4)
    bh = Inches(1.9)

    _plain_line(slide, bx, by, bx + bw, by, color=NAVY, width=2.0)
    _plain_line(slide, bx, by + bh, bx + bw, by + bh, color=NAVY, width=2.0)
    _plain_line(slide, bx, by, bx, by + bh, color=NAVY, width=2.0)
    _plain_line(slide, bx + bw, by, bx + bw, by + bh, color=NAVY, width=2.0)

    mid_y = by + Emu(int(bh / 2))
    _plain_line(slide, bx, mid_y, bx + Inches(2.4), mid_y, color=NAVY, width=2.0)
    _plain_line(slide, bx + bw - Inches(2.4), mid_y, bx + bw, mid_y, color=NAVY,
                width=2.0)

    _box(slide, bx + Inches(2.4), mid_y - Inches(0.35), Inches(1.6), Inches(0.7),
         "MOTOR", fill=WHITE, edge=NAVY, size=13, bold=True,
         shape=MSO_SHAPE.OVAL)

    _label(slide, bx - Inches(1.05), mid_y - Inches(0.5), Inches(1.0),
           "IN1\n(pin A)", size=12, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)
    _label(slide, bx + bw + Inches(0.1), mid_y - Inches(0.5), Inches(1.1),
           "IN2\n(pin B)", size=12, bold=True, color=TEAL)
    _label(slide, bx + Inches(2.6), by - Inches(0.36), Inches(1.4), "+ V",
           size=12, bold=True, color=GREY, align=PP_ALIGN.CENTER)
    _label(slide, bx + Inches(2.6), by + bh + Inches(0.06), Inches(1.4), "GND",
           size=12, bold=True, color=GREY, align=PP_ALIGN.CENTER)

    return slide


# ===================================================================
# 5. the 32-LED loop
# ===================================================================

def led_map(deck, title="Where every LED number is on the vehicle", speaker=None):
    slide = deck.blank(title, speaker=speaker or [
        "Put a vehicle on the desk in the same orientation as the diagram before you say anything. Students who cannot match the picture to the object will guess for the rest of the lesson.",
        "The numbers run in the order the data flows, not in any order that looks tidy from outside. Say so - it is the answer to the question they are about to ask.",
        "Get them to point at LED 0 and LED 31 on a real vehicle before moving on.",
        "This map is worth printing. They will use it in every lighting exercise from here to the end of the course.",
    ])

    left = MARGIN_L + Inches(0.9)
    width = Inches(10.6)
    led_w = Emu(int(width / 16))
    led_h = Inches(0.42)

    front_y = BODY_TOP + Inches(0.75)
    rear_y = BODY_TOP + Inches(2.85)

    _label(slide, left, front_y - Inches(0.4), width, "FRONT   -   data goes in at LED 0",
           size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    _label(slide, left, rear_y + led_h + Inches(0.1), width,
           "REAR   -   numbering runs the other way, right to left",
           size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    for i in range(16):
        fill = LIGHT_AMBER if i < 8 else LIGHT_TEAL
        _box(slide, left + led_w * i, front_y, led_w, led_h, str(i), fill=fill,
             edge=NAVY, size=11, shape=MSO_SHAPE.RECTANGLE, edge_w=0.75)

    for slot in range(16):
        number = 16 + slot                       # 16 at the right, 31 at the left
        fill = LIGHT_TEAL if number <= 23 else LIGHT_AMBER
        x = left + width - led_w * (slot + 1)    # so 16 sits under 15
        _box(slide, x, rear_y, led_w, led_h, str(number), fill=fill, edge=NAVY,
             size=11, shape=MSO_SHAPE.RECTANGLE, edge_w=0.75)

    # Side labels
    _label(slide, MARGIN_L - Inches(0.05), front_y + Inches(0.85), Inches(0.9),
           ["LEFT", "side"], size=13, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    _label(slide, left + width + Inches(0.06), front_y + Inches(0.85), Inches(1.0),
           ["RIGHT", "side"], size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # The mirror pairs
    for p in (0, 5, 10, 15):
        x = left + led_w * p + Emu(int(led_w / 2))
        _plain_line(slide, x, front_y + led_h, x, rear_y, color=RULE, width=1.0)

    _label(slide, left, front_y + Inches(1.05), width,
           "front LED p  and  rear LED 31 - p  are directly opposite each other",
           size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    _label(slide, left, front_y + Inches(1.38), width,
           "0 pairs with 31      5 pairs with 26      10 pairs with 21      15 pairs with 16",
           size=12, color=GREY, align=PP_ALIGN.CENTER, font=CODE_FONT)

    deck._note(slide,
               "LEFT is 0-7 and 24-31.  RIGHT is 8-15 and 16-23.  The strip is one "
               "loop, so each side is in two pieces.", "info")
    return slide


# ===================================================================
# 6. deadzone and map
# ===================================================================

def deadzone_map(deck, title="From thumbstick to motor speed", stick_max=127,
                 deadzone=20, speaker=None):
    slide = deck.blank(title, speaker=speaker or [
        "The two numbers that matter are the deadzone and the top of the stick range. Everything else on the slide follows from those.",
        "Ask why the line does not start at zero. The answer - a stick at rest is never exactly centred - is worth getting from the room rather than telling them.",
        "The jump at the deadzone edge is deliberate: below the minimum duty the motor buzzes instead of turning, so the map skips straight past it.",
        "Come back to this diagram when somebody's vehicle creeps on its own.",
    ])

    ox = MARGIN_L + Inches(1.3)
    oy = BODY_TOP + Inches(4.05)
    w = Inches(6.4)
    h = Inches(3.6)

    # Axes
    _plain_line(slide, ox, oy, ox + w, oy, color=NAVY, width=1.75)
    _plain_line(slide, ox, oy, ox, oy - h, color=NAVY, width=1.75)

    _label(slide, ox, oy + Inches(0.12), w,
           "stick reading   0  to  {}".format(stick_max), size=12, color=GREY,
           align=PP_ALIGN.CENTER)
    _label(slide, ox - Inches(1.25), oy - h - Inches(0.05), Inches(1.15),
           ["motor", "speed", "0 to 255"], size=12, color=GREY,
           align=PP_ALIGN.RIGHT)

    dz_frac = deadzone / float(stick_max)
    dz_x = ox + Emu(int(w * dz_frac))

    # Deadzone band
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, ox, oy - h,
                                  Emu(int(w * dz_frac)), h)
    band.fill.solid()
    band.fill.fore_color.rgb = LIGHT_RED
    band.line.fill.background()
    band.shadow.inherit = False
    band.text_frame.text = ""

    # Flat part, then the ramp
    _plain_line(slide, ox, oy, dz_x, oy, color=RED, width=3.0)
    _plain_line(slide, dz_x, oy, ox + w, oy - h, color=TEAL, width=3.0)

    _label(slide, ox, oy - h - Inches(0.42), Emu(int(w * dz_frac)) + Inches(0.6),
           "DEADZONE\nanswer is 0", size=11, bold=True, color=RED,
           align=PP_ALIGN.CENTER)

    _label(slide, dz_x + Inches(0.3), oy - h + Inches(0.1), Inches(4.0),
           "map() stretches what is left\nacross the whole speed range",
           size=13, bold=True, color=TEAL)

    _label(slide, dz_x - Inches(0.5), oy + Inches(0.36), Inches(1.2),
           str(deadzone), size=12, color=RED, align=PP_ALIGN.CENTER,
           font=CODE_FONT)

    # The code, alongside
    code_left = ox + w + Inches(1.15)
    code_w = Inches(3.6)
    _label(slide, code_left, BODY_TOP + Inches(0.15), code_w,
           "stickToSpeed()", size=14, bold=True, color=TEAL, font=CODE_FONT)

    lines = [
        "if (abs(v) < DEADZONE)",
        "    return 0;",
        "",
        "size = map(abs(v),",
        "           DEADZONE, {},".format(stick_max),
        "           MOTOR_MIN, maxSpeed);",
        "",
        "return (v > 0) ? size",
        "                : -size;",
    ]
    _box(slide, code_left, BODY_TOP + Inches(0.55), code_w, Inches(2.5), lines,
         fill=CODE_BG, edge=RULE, size=11, font=CODE_FONT,
         shape=MSO_SHAPE.RECTANGLE, edge_w=0.75,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    _label(slide, code_left, BODY_TOP + Inches(3.25), code_w,
           ["Why a deadzone at all?",
            "A worn stick does not return",
            "to exactly zero. Without this",
            "the vehicle creeps away on",
            "its own."],
           size=12, color=INK)

    return slide


# ===================================================================
# 7. tank drive mixing
# ===================================================================

def tank_mixing(deck, title="Mixing: one stick, two sides", speaker=None):
    slide = deck.blank(title, speaker=speaker or [
        "Two numbers in, two numbers out. Write the two lines on the board and work one example with the room: stick half forward and half right, what does each side get?",
        "The constrain() is not decoration. Without it, forward plus turn can ask for more than the motor has, and the vehicle turns the wrong way.",
        "Ask what happens with the stick pushed fully sideways: one side forward, one side back, and the vehicle spins on the spot. That is the move they need for the square.",
    ])

    _label(slide, MARGIN_L, BODY_TOP, CONTENT_W,
           "left = forward + turn            right = forward - turn",
           size=22, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    cases = [
        ("stick straight up", "forward 200\nturn 0", "left 200\nright 200",
         "straight ahead", LIGHT_TEAL),
        ("up and right", "forward 200\nturn 60", "left 260 -> 255\nright 140",
         "curves right\nwhile moving", LIGHT_TEAL),
        ("hard right only", "forward 0\nturn 127", "left 127\nright -127",
         "spins on the spot", LIGHT_AMBER),
        ("stick centred", "forward 0\nturn 0", "left 0\nright 0",
         "both sides coast", LIGHT_GREY),
    ]

    col_w = Inches(2.85)
    gap = Inches(0.28)
    top = BODY_TOP + Inches(0.75)

    for index, (what, inputs, outputs, meaning, fill) in enumerate(cases):
        left = MARGIN_L + (col_w + gap) * index

        _box(slide, left, top, col_w, Inches(0.45), what, fill=fill, edge=NAVY,
             size=13, bold=True, color=NAVY)

        _box(slide, left, top + Inches(0.6), col_w, Inches(0.85),
             inputs.split("\n"), fill=WHITE, edge=RULE, size=13, font=CODE_FONT,
             shape=MSO_SHAPE.RECTANGLE, edge_w=0.75)

        _arrow(slide, left + Emu(int(col_w / 2)), top + Inches(1.5),
               left + Emu(int(col_w / 2)), top + Inches(1.85))

        _box(slide, left, top + Inches(1.9), col_w, Inches(0.85),
             outputs.split("\n"), fill=CODE_BG, edge=TEAL, size=13,
             font=CODE_FONT, shape=MSO_SHAPE.RECTANGLE, edge_w=1.0)

        _label(slide, left, top + Inches(2.9), col_w, meaning.split("\n"),
               size=13, color=INK, align=PP_ALIGN.CENTER)

    deck._note(slide,
               "260 is over the maximum, which is exactly what constrain() is "
               "there to catch. Steering also has its own lower cap, turnMax, "
               "so the vehicle stays easy to aim.", "info")
    return slide


# ===================================================================
# 8. delay vs millis
# ===================================================================

def millis_timeline(deck, title="Why the vehicle programs never call delay()", speaker=None):
    slide = deck.blank(title, speaker=speaker or [
        "This is the slide that explains the whole of Lesson 5. Do not rush it.",
        "Two jobs at different rates cannot both be done with delay(). Draw the two timelines on the board and let them see the collision.",
        "millis() is not a timer you start. It is a clock that has been running since power-up, and all you ever do is subtract.",
        "Subtracting rather than comparing is what survives the wrap at 49 days. Mention it, do not dwell on it.",
    ])

    left = MARGIN_L + Inches(1.9)
    width = Inches(9.3)

    # ---- with delay ----
    y = BODY_TOP + Inches(0.45)
    _label(slide, MARGIN_L, y - Inches(0.48), Inches(7.0),
           "with delay()  -  one job, and nothing else runs",
           size=13, bold=True, color=RED)

    segments = [("blink", 0.10, LIGHT_TEAL), ("frozen - delay(1000)", 0.40, LIGHT_RED),
                ("blink", 0.10, LIGHT_TEAL), ("frozen - delay(1000)", 0.40, LIGHT_RED)]
    x = left
    for name, frac, fill in segments:
        seg_w = Emu(int(width * frac))
        _box(slide, x, y, seg_w, Inches(0.55), name, fill=fill, edge=RULE,
             size=11, shape=MSO_SHAPE.RECTANGLE, edge_w=0.75)
        x += seg_w

    _label(slide, left, y + Inches(0.62), width,
           "while it is frozen it is not reading the controller, not checking the "
           "battery, and not updating any other light",
           size=12, color=RED, align=PP_ALIGN.CENTER)

    # ---- with millis ----
    # The heading goes ABOVE the rows rather than beside them: at the 18pt
    # floor a two-line label in the left gutter lands on top of the job names.
    y = BODY_TOP + Inches(2.05)
    _label(slide, MARGIN_L, y, Inches(7.0),
           "with millis()  -  many jobs, none of them waiting",
           size=13, bold=True, color=TEAL)
    y += Inches(0.50)

    jobs = [
        ("drive motors", 0.02, TEAL),
        ("scanner frame", 0.08, AMBER),
        ("turn signal", 0.16, NAVY),
        ("battery check", 0.5, GREY),
    ]
    row_h = Inches(0.40)
    for row, (name, spacing, colour) in enumerate(jobs):
        ry = y + row_h * row
        _label(slide, left - Inches(1.85), ry + Inches(0.04), Inches(1.75), name,
               size=11, color=colour, align=PP_ALIGN.RIGHT)
        _plain_line(slide, left, ry + Inches(0.18), left + width, ry + Inches(0.18),
                    color=RULE, width=0.75)
        step = max(spacing, 0.02)
        count = int(1.0 / step)
        for tick in range(count + 1):
            tx = left + Emu(int(width * step * tick))
            if tx > left + width:
                break
            _plain_line(slide, tx, ry + Inches(0.06), tx, ry + Inches(0.30),
                        color=colour, width=2.0)

    _label(slide, left, y + row_h * 4 + Inches(0.12), width,
           "each job fires on its own schedule, and loop() never stops running",
           size=12, color=TEAL, align=PP_ALIGN.CENTER)

    _box(slide, MARGIN_L, BODY_TOP + Inches(4.55), CONTENT_W, Inches(0.72),
         "if (millis() - lastTime >= INTERVAL) {  lastTime = millis();  ...do the thing...  }",
         fill=CODE_BG, edge=TEAL, size=15, font=CODE_FONT,
         shape=MSO_SHAPE.RECTANGLE, edge_w=1.0)

    return slide


# ===================================================================
# 9. dead reckoning square
# ===================================================================

def square_path(deck, title="Dead reckoning: the vehicle has no idea where it is", speaker=None):
    """
    The square the vehicle drives, with the four legs numbered and the two
    equations that predict it. Redrawn so the corner labels sit outside the
    path instead of on top of it.
    """
    slide = deck.blank(title, speaker=speaker or [
        "The vehicle has no idea where it is. Say that sentence out loud - it is the point of the diagram and it surprises people.",
        "Everything it does here is open loop: a speed, for a time, and then a hope.",
        "Ask what makes the real path miss the drawn one. Battery voltage, floor surface, one motor slightly faster than another, wheel slip on the turn.",
        "This is the honest introduction to why sensors exist, and it sets up the advanced course.",
    ])

    side = Inches(2.9)
    ox = MARGIN_L + Inches(1.5)
    oy = BODY_TOP + Inches(3.60)

    corners = [(ox, oy), (ox, oy - side), (ox + side, oy - side), (ox + side, oy)]

    # The four legs, drawn anticlockwise from the start.
    for index in range(4):
        x1, y1 = corners[index]
        x2, y2 = corners[(index + 1) % 4]
        _arrow(slide, x1, y1, x2, y2, color=TEAL, width=3.0)

    # Number each leg on the OUTSIDE of the path.
    leg_labels = [
        (ox - Inches(1.75), oy - Emu(int(side / 2)) - Inches(0.2), "1  forward"),
        (ox + Emu(int(side / 2)) - Inches(0.8), oy - side - Inches(0.75), "2  forward"),
        (ox + side + Inches(0.3), oy - Emu(int(side / 2)) - Inches(0.2), "3  forward"),
        (ox + Emu(int(side / 2)) - Inches(0.8), oy + Inches(0.35), "4  forward"),
    ]
    for lx, ly, text in leg_labels:
        _label(slide, lx, ly, Inches(1.6), text, size=14, bold=True,
               color=NAVY, align=PP_ALIGN.CENTER)

    # Corner markers, small, just outside each turn.
    for index, (x, y) in enumerate(corners):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x - Inches(0.13),
                                     y - Inches(0.13), Inches(0.26), Inches(0.26))
        dot.fill.solid()
        dot.fill.fore_color.rgb = AMBER
        dot.line.fill.background()
        dot.shadow.inherit = False
        _set_text(dot.text_frame, [""], size=8)

    _label(slide, ox - Inches(1.95), oy - Inches(0.30), Inches(1.85),
           "START", size=15, bold=True, color=TEAL, align=PP_ALIGN.RIGHT)
    _label(slide, ox - Inches(1.95), oy + Inches(0.06), Inches(1.85),
           "each dot is a", size=13, color=AMBER, align=PP_ALIGN.RIGHT)
    _label(slide, ox - Inches(1.95), oy + Inches(0.42), Inches(1.85),
           "90 degree turn", size=13, color=AMBER, align=PP_ALIGN.RIGHT)

    # The maths, well clear of the figure.
    right = MARGIN_L + Inches(6.6)
    col = Inches(5.6)

    _label(slide, right, BODY_TOP, col, "The two equations", size=19,
           bold=True, color=TEAL)

    _box(slide, right, BODY_TOP + Inches(0.46), col, Inches(0.98),
         ["distance = speed x time",
          "degrees turned = turn rate x time"],
         fill=CODE_BG, edge=TEAL, size=15, font=CODE_FONT,
         shape=MSO_SHAPE.RECTANGLE, edge_w=1.0, align=PP_ALIGN.LEFT)

    _label(slide, right, BODY_TOP + Inches(1.56), col,
           ["You know the TIME - you wrote it in the delay().",
            "You have to MEASURE the speed and the turn rate.",
            "wheel circumference = pi x 3.0 in = 9.42 in",
            "one wheel turn = 9.42 / 12 = 0.79 ft",
            "",
            "Drive one side, time it, divide: that is feet per",
            "second for YOUR vehicle, on THIS floor and battery."],
           size=14, color=INK)

    deck._note(slide,
               "Recharge the battery and the same numbers give a bigger square. "
               "That drift is why dead reckoning is a starting point, not an "
               "answer.", "warn")
    return slide


# ===================================================================
# 10. RGB additive mixing
# ===================================================================

def rgb_mixing(deck, title="One pixel is three LEDs, and colour is a mixture", speaker=None):
    slide = deck.blank(title, speaker=speaker or [
        "One pixel is three tiny LEDs behind one lens. Hold a vehicle close enough that somebody in the front row can see them.",
        "Three numbers, 0 to 255 each, is 16.7 million colours. Worth doing the multiplication on the board.",
        "Warn them now: 255, 255, 255 is every LED at full, and it is both blinding and expensive in current. The power budget slide follows for a reason.",
        "If time allows, get them to predict a colour before they type it, then check. Predicting is what makes it stick.",
    ])

    cx = MARGIN_L + Inches(3.0)
    cy = BODY_TOP + Inches(2.35)
    r = Inches(1.5)

    circles = [
        ("R", RGBColor(0xE8, 0x2C, 0x2C), cx, cy - Inches(0.85)),
        ("G", RGBColor(0x2C, 0xB4, 0x4A), cx - Inches(0.85), cy + Inches(0.6)),
        ("B", RGBColor(0x2C, 0x5A, 0xE8), cx + Inches(0.85), cy + Inches(0.6)),
    ]

    for name, colour, x, y in circles:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x - r, y - r, r * 2, r * 2)
        circle.fill.solid()
        circle.fill.fore_color.rgb = colour
        circle.fill.transparency = 0.45
        circle.line.color.rgb = colour
        circle.line.width = Pt(1.5)
        circle.shadow.inherit = False
        _set_text(circle.text_frame, [""], size=8)

    _label(slide, MARGIN_L, cy + Inches(1.86), Inches(6.0),
           "where all three overlap you get white", size=13, color=GREY,
           align=PP_ALIGN.CENTER)

    right = MARGIN_L + Inches(6.6)
    col_w = Inches(6.1)

    _label(slide, right, BODY_TOP + Inches(0.05), col_w,
           "strip.Color(red, green, blue)", size=17, bold=True, color=TEAL,
           font=CODE_FONT)
    _label(slide, right, BODY_TOP + Inches(0.45), col_w,
           "Each number is 0 to 255 - how hard that one LED is driven.",
           size=14, color=INK)

    swatches = [
        ("(255,   0,   0)", "red", RGBColor(0xE8, 0x2C, 0x2C)),
        ("(  0, 255,   0)", "green", RGBColor(0x2C, 0xB4, 0x4A)),
        ("(  0,   0, 255)", "blue", RGBColor(0x2C, 0x5A, 0xE8)),
        ("(255, 255,   0)", "yellow", RGBColor(0xE8, 0xC8, 0x2C)),
        ("(255, 100,   0)", "amber - the turn signals", RGBColor(0xE8, 0x8A, 0x1A)),
        ("(255, 255, 255)", "white - and three times the current",
         RGBColor(0xDD, 0xDD, 0xDD)),
    ]

    row_h = Inches(0.52)
    for index, (code, name, colour) in enumerate(swatches):
        y = BODY_TOP + Inches(1.05) + row_h * index
        chip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right, y,
                                      Inches(0.55), Inches(0.34))
        chip.fill.solid()
        chip.fill.fore_color.rgb = colour
        chip.line.color.rgb = GREY
        chip.line.width = Pt(0.75)
        chip.shadow.inherit = False
        _set_text(chip.text_frame, [""], size=8)

        _label(slide, right + Inches(0.75), y + Inches(0.02), Inches(1.9), code,
               size=13, font=CODE_FONT, color=INK)
        _label(slide, right + Inches(2.8), y + Inches(0.02), Inches(3.3), name,
               size=13, color=GREY)

    deck._note(slide,
               "Each of the three LEDs draws about 20 mA. White costs three times "
               "what red costs, and 32 pixels of white is nearly 2 amps.", "warn")
    return slide


# ===================================================================
# 11. servo pulse widths
# ===================================================================

def servo_pulse(deck, title="A servo listens to the LENGTH of a pulse", speaker=None):
    slide = deck.blank(title, speaker=speaker or [
        "A servo does not listen to voltage, or to duty in the way a motor does. It listens to how LONG the pulse is, and it wants one every 20 ms.",
        "1.0 ms is one end, 1.5 ms is centre, 2.0 ms is the other end. Those three numbers are the whole protocol.",
        "The ESP32 counts duty, not microseconds, so the program has to convert. Walk the arithmetic once, slowly.",
        "Continuous-rotation servos read the same pulse as a SPEED rather than a position. Worth flagging before somebody buys the wrong one.",
    ])

    left = MARGIN_L + Inches(1.7)
    width = Inches(8.0)
    top = BODY_TOP + Inches(0.18)

    rows = [
        ("1.0 ms", 0.05, "one end of the travel", TEAL),
        ("1.5 ms", 0.075, "centred", NAVY),
        ("2.0 ms", 0.10, "the other end", AMBER),
    ]

    for index, (name, frac, meaning, colour) in enumerate(rows):
        y = top + Inches(1.12) * index
        base = y + Inches(0.68)

        _label(slide, MARGIN_L, y + Inches(0.2), Inches(1.6), name, size=15,
               bold=True, color=colour, font=CODE_FONT)

        _plain_line(slide, left, base, left + width, base, color=RULE, width=1.0)

        pulse_w = Emu(int(width * frac))
        _plain_line(slide, left, base, left, y + Inches(0.1), color=colour, width=2.5)
        _plain_line(slide, left, y + Inches(0.1), left + pulse_w, y + Inches(0.1),
                    color=colour, width=2.5)
        _plain_line(slide, left + pulse_w, y + Inches(0.1), left + pulse_w, base,
                    color=colour, width=2.5)

        _label(slide, left + pulse_w + Inches(0.15), y + Inches(0.16), Inches(4.5),
               meaning, size=13, color=INK)

    _label(slide, left, top + Inches(3.46), width,
           "the whole cycle is 20 ms, repeated 50 times a second",
           size=12, color=GREY, align=PP_ALIGN.CENTER)

    _box(slide, MARGIN_L, BODY_TOP + Inches(4.18), CONTENT_W, Inches(0.85),
         ["The ESP32 counts duty, not microseconds. At 16-bit resolution one 20 ms cycle is 65536 counts:",
          "duty = microseconds * 65536 / 20000"],
         fill=CODE_BG, edge=TEAL, size=14, font=CODE_FONT,
         shape=MSO_SHAPE.RECTANGLE, edge_w=1.0, align=PP_ALIGN.LEFT)

    return slide


# ===================================================================
# 12. lighting state machine
# ===================================================================

def state_machine(deck, title="The lighting state machine", speaker=None):
    slide = deck.blank(title, speaker=speaker or [
        "A state machine is just: what am I doing now, and what could make that change. The diagram is a picture of exactly that.",
        "Point out that exactly one state is active at a time. That is what stops the lights contradicting themselves.",
        "Ask the room for a transition that is missing, then decide together whether it should be there.",
        "This shape turns up everywhere in embedded work. Naming it now means they recognise it later.",
    ])

    top = BODY_TOP + Inches(0.35)

    nodes = {
        "DISCONNECTED": (MARGIN_L + Inches(0.2), top, LIGHT_GREY, "green pulse\nno controller"),
        "PAIRING": (MARGIN_L + Inches(0.2), top + Inches(2.1), LIGHT_TEAL, "blue breathing\nwaiting to pair"),
        "STANDBY": (MARGIN_L + Inches(4.6), top + Inches(1.05), WHITE, "headlights + tail lights\nredrawn only when dirty"),
        "TURN_LEFT": (MARGIN_L + Inches(9.0), top - Inches(0.25), LIGHT_AMBER, "amber bar, 3 cycles"),
        "TURN_RIGHT": (MARGIN_L + Inches(9.0), top + Inches(1.05), LIGHT_AMBER, "amber bar, 3 cycles"),
        "KITT_SCANNER": (MARGIN_L + Inches(9.0), top + Inches(2.35), LIGHT_TEAL, "sweeping dot, 3 cycles"),
    }

    node_w = Inches(3.1)
    node_h = Inches(0.62)
    placed = {}

    for name, (x, y, fill, blurb) in nodes.items():
        _box(slide, x, y, node_w, node_h, name, fill=fill, edge=NAVY, size=13,
             bold=True, color=NAVY, font=CODE_FONT)
        _label(slide, x, y + node_h + Inches(0.03), node_w, blurb.split("\n"),
               size=11, color=GREY, align=PP_ALIGN.CENTER)
        placed[name] = (x, y)

    def centre_right(name):
        x, y = placed[name]
        return x + node_w, y + Emu(int(node_h / 2))

    def centre_left(name):
        x, y = placed[name]
        return x, y + Emu(int(node_h / 2))

    for source in ("DISCONNECTED", "PAIRING"):
        x1, y1 = centre_right(source)
        x2, y2 = centre_left("STANDBY")
        _arrow(slide, x1, y1, x2, y2, color=TEAL)

    for target in ("TURN_LEFT", "TURN_RIGHT", "KITT_SCANNER"):
        x1, y1 = centre_right("STANDBY")
        x2, y2 = centre_left(target)
        _arrow(slide, x1, y1, x2, y2, color=AMBER)

    _label(slide, MARGIN_L + Inches(3.5), top + Inches(0.55), Inches(1.4),
           "controller\nconnects", size=10, color=TEAL, align=PP_ALIGN.CENTER)
    _label(slide, MARGIN_L + Inches(7.75), top + Inches(0.35), Inches(1.3),
           "D-pad\nor button", size=10, color=AMBER, align=PP_ALIGN.CENTER)
    _label(slide, MARGIN_L + Inches(7.6), top + Inches(2.75), Inches(1.5),
           "each animation counts\nits own cycles and\nhands itself back",
           size=10, color=GREY, align=PP_ALIGN.CENTER)

    deck._note(slide,
               "One variable holds the mode, so two modes can never be half-on at "
               "once. Adding a mode is one enum value and one branch - not a new "
               "combination of flags to test.", "info")
    return slide


# ===================================================================
# 13. current signature
# ===================================================================

def current_signature(deck, title="What a healthy motor looks like to a current sensor", speaker=None):
    slide = deck.blank(title, speaker=speaker or [
        "Current is the cheapest sensor on the vehicle, because it is already there. Every fault has a shape.",
        "Walk the three regions: no load, working, stalled. A stalled motor draws the most and moves the least, which is the worst of both.",
        "Ask what a broken wire looks like. Near zero current with the command still on - and the vehicle limping in a circle.",
        "This is the foundation of the powered self-test later in the lesson.",
    ])

    left = MARGIN_L + Inches(1.15)
    width = Inches(3.3)
    gap = Inches(0.5)
    top = BODY_TOP + Inches(0.50)
    height = Inches(2.1)

    traces = [
        ("HEALTHY", [0.08, 0.08, 0.95, 0.62, 0.45, 0.40, 0.38, 0.38],
         "spikes as it breaks away,\nthen settles", TEAL, LIGHT_TEAL),
        ("DISCONNECTED", [0.06, 0.06, 0.08, 0.07, 0.06, 0.07, 0.06, 0.06],
         "no spike, almost no draw", GREY, LIGHT_GREY),
        ("JAMMED", [0.08, 0.08, 0.98, 0.95, 0.94, 0.96, 0.95, 0.94],
         "draws a lot and never settles", RED, LIGHT_RED),
    ]

    for index, (name, points, meaning, colour, fill) in enumerate(traces):
        x0 = left + (width + gap) * index
        base = top + height

        panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, top, width, height)
        panel.fill.solid()
        panel.fill.fore_color.rgb = fill
        panel.line.color.rgb = RULE
        panel.line.width = Pt(0.75)
        panel.shadow.inherit = False
        _set_text(panel.text_frame, [""], size=8)

        # baseline marker
        base_y = base - Emu(int(height * 0.10))
        _plain_line(slide, x0, base_y, x0 + width, base_y, color=GREY, width=1.0)

        step = Emu(int(width / (len(points) - 1)))
        for i in range(len(points) - 1):
            y1 = base - Emu(int(height * points[i] * 0.92))
            y2 = base - Emu(int(height * points[i + 1] * 0.92))
            _plain_line(slide, x0 + step * i, y1, x0 + step * (i + 1), y2,
                        color=colour, width=2.5)

        _label(slide, x0, top - Inches(0.38), width, name, size=15, bold=True,
               color=colour, align=PP_ALIGN.CENTER)
        _label(slide, x0, base + Inches(0.1), width, meaning.split("\n"), size=12,
               color=INK, align=PP_ALIGN.CENTER)

    _label(slide, left - Inches(1.15), top + height - Inches(0.34),
           Inches(1.05), "baseline",
           size=12, color=GREY, align=PP_ALIGN.RIGHT)

    _box(slide, MARGIN_L, BODY_TOP + Inches(3.40), CONTENT_W, Inches(0.80),
         ["Both thresholds are measured ABOVE the baseline that was just taken, not as absolute amps:",
          "peak - baseline >= 1.500 A     AND     average - baseline >= 1.000 A"],
         fill=CODE_BG, edge=TEAL, size=14, font=CODE_FONT,
         shape=MSO_SHAPE.RECTANGLE, edge_w=1.0, align=PP_ALIGN.LEFT)

    deck._note(slide,
               "Op 11.2 measured a baseline, printed it, then compared against "
               "fixed absolutes anyway - so a vehicle idling at half an amp failed "
               "every motor. Op 12 fixed it.", "warn")
    return slide


# ===================================================================
# 14. Ohm's law and the power law
# ===================================================================

def ohms_and_power_law(deck,
                       title="Ohm's law, and the power law you need for the LEDs", speaker=None):
    """
    The two relationships the course keeps coming back to: sizing the LED
    current budget, reading the shunt resistor, and working out battery life.
    Drawn rather than borrowed, so it stays editable and prints sharp.
    """
    slide = deck.blank(title, speaker=speaker or [
        "Two relationships, and between them they explain every resistor and every current budget on this vehicle.",
        "Do one worked example on the board before they look at the slide, and make them tell you which quantity is the unknown.",
        "Units matter more than the algebra here. Volts, amps, ohms, watts - keep saying them.",
        "The reason this is in a robotics course: the LED resistor and the 32-LED current budget are both this, and getting them wrong smells of hot plastic.",
    ])

    top = BODY_TOP + Inches(0.35)
    tri_w = Inches(3.0)
    tri_h = Inches(2.2)

    groups = [
        ("OHM'S LAW", MARGIN_L + Inches(0.7), "V", "I", "R", TEAL, LIGHT_TEAL,
         ["V = I x R", "I = V / R", "R = V / I"],
         ["V  volts", "I  amps", "R  ohms"]),
        ("POWER LAW", MARGIN_L + Inches(6.9), "P", "I", "V", AMBER, LIGHT_AMBER,
         ["P = I x V", "I = P / V", "V = P / I"],
         ["P  watts", "I  amps", "V  volts"]),
    ]

    for name, left, apex, bl, br, edge, fill, formulas, legend in groups:
        _label(slide, left, top - Inches(0.42), tri_w, name, size=16, bold=True,
               color=edge, align=PP_ALIGN.CENTER)

        tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top,
                                     tri_w, tri_h)
        tri.fill.solid()
        tri.fill.fore_color.rgb = fill
        tri.line.color.rgb = edge
        tri.line.width = Pt(1.5)
        tri.shadow.inherit = False
        _set_fitted(tri.text_frame, [""], width=tri_w, height=tri_h, size=10)

        # The divider under the apex, the way the mnemonic is always drawn.
        mid_y = top + Emu(int(tri_h * 0.56))
        _plain_line(slide, left + Inches(0.62), mid_y,
                    left + tri_w - Inches(0.62), mid_y, color=edge, width=1.75)

        _label(slide, left, top + Inches(0.42), tri_w, apex, size=26, bold=True,
               color=NAVY, align=PP_ALIGN.CENTER)
        _label(slide, left + Inches(0.35), mid_y + Inches(0.30), Inches(0.9), bl,
               size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _label(slide, left + tri_w - Inches(1.25), mid_y + Inches(0.30),
               Inches(0.9), br, size=22, bold=True, color=NAVY,
               align=PP_ALIGN.CENTER)

        _box(slide, left - Inches(0.15), top + tri_h + Inches(0.3),
             tri_w + Inches(0.3), Inches(1.0), formulas, fill=CODE_BG,
             edge=RULE, size=14, font=CODE_FONT, shape=MSO_SHAPE.RECTANGLE,
             edge_w=0.75, align=PP_ALIGN.CENTER)

        _label(slide, left - Inches(0.15), top + tri_h + Inches(1.42),
               tri_w + Inches(0.3), "     ".join(legend), size=12, color=GREY,
               align=PP_ALIGN.CENTER)

    # Worked example, between the two triangles.
    mid_left = MARGIN_L + Inches(4.15)
    _box(slide, mid_left, top - Inches(0.2), Inches(2.4), Inches(4.0),
         ["ONE NEOPIXEL",
          "at full white",
          "",
          "3 LEDs x 20 mA",
          "= 60 mA",
          "",
          "P = I x V",
          "= 0.06 x 5",
          "= 0.3 W",
          "",
          "x 32 pixels",
          "= 1.92 A",
          "= 9.6 W"],
         fill=PANEL, edge=TEAL, size=13, shape=MSO_SHAPE.RECTANGLE,
         edge_w=1.0, align=PP_ALIGN.CENTER)

    deck._note(slide,
               "Cover the quantity you want with your thumb and the triangle "
               "shows you the sum. These two turn up again in the LED budget, "
               "the shunt resistor, and how long your battery lasts.", "info")
    return slide


# ===================================================================
# 15. the LED circuit students wire in Lesson 1
# ===================================================================

def led_circuit(deck, title="The circuit you are about to build", speaker=None):
    """
    GPIO -> resistor -> LED -> ground, drawn as a loop, with the Ohm's law
    working that picks the resistor. Native shapes, so it prints sharp and a
    teacher can retype the numbers for a different LED.
    """
    slide = deck.blank(title, speaker=speaker or [
        "Trace the loop with a finger: out of the pin, through the resistor, through the LED, back to ground. A circuit that is not a loop does nothing.",
        "Ask what the resistor is for BEFORE you say. The answer is that the LED will take as much current as you let it, and then stop being an LED.",
        "The long leg is the anode and goes towards the pin. Backwards means no light and no damage, so let them find out.",
        "Ground is not optional. Half the circuits that do not work are missing the return path.",
    ])

    left = MARGIN_L + Inches(0.6)
    top = BODY_TOP + Inches(0.55)
    w = Inches(6.6)
    h = Inches(2.6)

    # The loop itself.
    _plain_line(slide, left, top, left + w, top, color=NAVY, width=2.5)
    _plain_line(slide, left + w, top, left + w, top + h, color=NAVY, width=2.5)
    _plain_line(slide, left, top + h, left + w, top + h, color=NAVY, width=2.5)
    _plain_line(slide, left, top, left, top + h, color=NAVY, width=2.5)

    _box(slide, left - Inches(1.0), top - Inches(0.32), Inches(2.0),
         Inches(0.64), "ESP32  GPIO 2", fill=LIGHT_TEAL, edge=TEAL, size=14,
         bold=True, color=NAVY)
    _box(slide, left - Inches(0.95), top + h - Inches(0.32), Inches(1.9),
         Inches(0.64), "ESP32  GND", fill=LIGHT_GREY, edge=GREY, size=14,
         bold=True, color=NAVY)

    _box(slide, left + Inches(1.6), top - Inches(0.34), Inches(2.1),
         Inches(0.68), "220 ohm", fill=WHITE, edge=AMBER, size=15, bold=True,
         color=NAVY)
    _box(slide, left + w - Inches(0.55), top + Inches(0.85), Inches(1.1),
         Inches(0.9), "LED", fill=LIGHT_AMBER, edge=AMBER, size=15, bold=True,
         color=NAVY)

    _label(slide, left + w - Inches(2.35), top + Inches(0.72), Inches(1.7),
           "long leg  +", size=14, color=TEAL, bold=True, align=PP_ALIGN.RIGHT)
    _label(slide, left + w - Inches(2.35), top + Inches(1.68), Inches(1.7),
           "short leg  -", size=14, color=GREY, align=PP_ALIGN.RIGHT)

    _label(slide, left, top + h + Inches(0.35), w,
           "Current goes round ONE loop. The resistor may sit on either side "
           "of the LED - what matters is that it is in the loop.",
           size=14, color=INK, align=PP_ALIGN.CENTER)

    # The sizing calculation.
    right = MARGIN_L + Inches(8.1)
    col = Inches(4.1)

    _label(slide, right, BODY_TOP + Inches(0.1), col, "Why 220 ohms?",
           size=18, bold=True, color=TEAL)

    _box(slide, right, BODY_TOP + Inches(0.6), col, Inches(1.65),
         ["R = (supply - LED drop) / current",
          "",
          "  = (3.3 - 2.0) / 0.020",
          "  = 65 ohms"],
         fill=CODE_BG, edge=RULE, size=14, font=CODE_FONT,
         shape=MSO_SHAPE.RECTANGLE, edge_w=0.75, align=PP_ALIGN.LEFT)

    _label(slide, right, BODY_TOP + Inches(2.4), col,
           ["220 ohms is in the kit, and is a good choice:",
            "",
            "(3.3 - 2.0) / 220 = about 6 mA",
            "",
            "Dimmer than flat out, and easy to see."],
           size=14, color=INK)

    deck._note(slide,
               "No resistor means one bright flash and a dead LED. An ESP32 pin "
               "should not be asked for more than about 20 mA, so erring high "
               "is the right way to err.", "warn")
    return slide


# ===================================================================
# 16. series and parallel
# ===================================================================

def series_parallel(deck, title="Series and parallel: two ways to wire two loads", speaker=None):
    """
    The two circuit shapes side by side with the rules that go with them.
    Redrawn rather than lifted, because Kevin's originals are screenshots of a
    PDF in a browser window.
    """
    slide = deck.blank(title, speaker=speaker or [
        "Same two parts, two different wirings, completely different behaviour. That is the whole slide.",
        "Series: one path, so the same current everywhere and the voltage shared out. Parallel: two paths, so the same voltage across each and the currents add.",
        "Ask which one the 32 LEDs must be, and why. The answer - parallel, because each needs the full supply voltage - leads straight into the current budget.",
        "Christmas lights are the story everyone already knows: one bulb out and the whole string dies means series.",
    ])

    gap = Inches(0.7)
    col = Emu(int((CONTENT_W - gap) / 2))
    top = BODY_TOP + Inches(0.5)

    for index, kind in enumerate(("series", "parallel")):
        left = MARGIN_L + (col + gap) * index
        edge = TEAL if kind == "series" else AMBER
        fill = LIGHT_TEAL if kind == "series" else LIGHT_AMBER

        _label(slide, left, BODY_TOP, col,
               "SERIES" if kind == "series" else "PARALLEL",
               size=18, bold=True, color=edge, align=PP_ALIGN.CENTER)

        bx, by = left + Inches(0.5), top + Inches(0.25)
        bw, bh = col - Inches(1.0), Inches(1.9)

        # Battery on the left edge of both.
        _plain_line(slide, bx, by, bx, by + bh, color=NAVY, width=2.0)
        _label(slide, bx - Inches(0.95), by + Inches(0.72), Inches(0.85),
               "battery", size=13, color=GREY, align=PP_ALIGN.RIGHT)

        if kind == "series":
            # One loop, two resistors in line.
            _plain_line(slide, bx, by, bx + bw, by, color=NAVY, width=2.0)
            _plain_line(slide, bx + bw, by, bx + bw, by + bh, color=NAVY, width=2.0)
            _plain_line(slide, bx, by + bh, bx + bw, by + bh, color=NAVY, width=2.0)
            _box(slide, bx + Inches(0.7), by - Inches(0.28), Inches(1.2),
                 Inches(0.56), "R1", fill=fill, edge=edge, size=14, bold=True)
            _box(slide, bx + bw - Inches(1.9), by - Inches(0.28), Inches(1.2),
                 Inches(0.56), "R2", fill=fill, edge=edge, size=14, bold=True)
            rules = ["One path. The SAME current everywhere.",
                     "",
                     "R total = R1 + R2",
                     "I total = I1 = I2",
                     "V total = V1 + V2"]
        else:
            # Two branches between the same two rails.
            mid = bx + Emu(int(bw / 2))
            _plain_line(slide, bx, by, bx + bw, by, color=NAVY, width=2.0)
            _plain_line(slide, bx, by + bh, bx + bw, by + bh, color=NAVY, width=2.0)
            _plain_line(slide, bx + bw, by, bx + bw, by + bh, color=NAVY, width=2.0)
            _plain_line(slide, mid, by, mid, by + bh, color=NAVY, width=2.0)
            _box(slide, mid - Inches(0.6), by + Inches(0.62), Inches(1.2),
                 Inches(0.6), "R1", fill=fill, edge=edge, size=14, bold=True)
            _box(slide, bx + bw - Inches(0.6), by + Inches(0.62), Inches(1.2),
                 Inches(0.6), "R2", fill=fill, edge=edge, size=14, bold=True)
            rules = ["Two paths. The SAME voltage across each.",
                     "",
                     "V total = V1 = V2",
                     "I total = I1 + I2",
                     "1/R total = 1/R1 + 1/R2"]

        _box(slide, left, top + Inches(2.5), col, Inches(1.75), rules,
             fill=WHITE, edge=RULE, size=15, shape=MSO_SHAPE.RECTANGLE,
             edge_w=0.75, align=PP_ALIGN.CENTER)

    deck._note(slide,
               "The LED circuit you built in Lesson 1 is a SERIES circuit: one "
               "loop, so the same current flows through the resistor and the "
               "LED. That is why the resistor protects the LED wherever you "
               "put it.", "info")
    return slide


# ===================================================================
# 17. six degrees of freedom
# ===================================================================

def six_dof(deck, title="Six degrees of freedom, and the two your rover has",
            controller="PS3 controller", speaker=None):
    """
    Surge, sway and heave against roll, pitch and yaw, drawn on a plan and a
    side view of the vehicle rather than on a ship.

    Redrawn from Kevin's Lesson 3 figure. His used a ship, which is the right
    picture for a submersible and a confusing one for somebody holding a
    rover, so the same six arrows are drawn on the vehicle in front of them.
    The point of the figure is the bottom row: a rover on a flat floor
    controls two of the six, which is why one thumbstick is enough.
    """
    slide = deck.blank(title, speaker=speaker or [
        "Six ways a rigid body can move. Three slides, three twists, and "
        "nothing else.",
        "Demonstrate with your hand held flat before you show the diagram. "
        "Slide it forward, sideways, up. Then tip it, nose it down, spin it.",
        "The two boxed in teal are the two a rover on a flat floor can "
        "actually control. That is why one stick with two axes is enough "
        "here.",
        "A submersible or an aircraft controls all six, which is why an ROV "
        "pilot needs two sticks and far more practice.",
        "Sway is the interesting one to ask about. A rover cannot slide "
        "sideways - it has to yaw and then surge. Ask them why.",
    ])

    top = BODY_TOP + Inches(0.30)
    col_w = Inches(5.9)
    right = MARGIN_L + col_w + Inches(0.55)

    _label(slide, MARGIN_L, top, col_w, "TRANSLATION  -  sliding along an axis",
           size=16, bold=True, color=TEAL)
    _label(slide, right, top, col_w, "ROTATION  -  turning about an axis",
           size=16, bold=True, color=TEAL)

    row_top = top + Inches(0.46)
    row_h = Inches(0.78)
    gap = Inches(0.14)

    # left column: the three translations
    translations = [
        ("SURGE", "forward and back", "X axis", True),
        ("SWAY", "left and right", "Y axis", False),
        ("HEAVE", "up and down", "Z axis", False),
    ]
    rotations = [
        ("ROLL", "tipping left or right", "about X", False),
        ("PITCH", "nose up or nose down", "about Y", False),
        ("YAW", "turning left or right", "about Z", True),
    ]

    for column, rows in ((MARGIN_L, translations), (right, rotations)):
        for index, (name, meaning, axis, ours) in enumerate(rows):
            box_top = row_top + (row_h + gap) * index
            fill = LIGHT_TEAL if ours else WHITE
            edge = TEAL if ours else RULE
            _box(slide, column, box_top, Inches(1.55), row_h, name,
                 fill=fill, edge=edge, size=17, bold=True, color=NAVY,
                 edge_w=2.0 if ours else 1.0)
            _label(slide, column + Inches(1.70), box_top + Inches(0.04),
                   Inches(2.55), meaning, size=15, height=Inches(0.36))
            _label(slide, column + Inches(4.35), box_top + Inches(0.10),
                   Inches(1.5), axis, size=14, color=GREY)
            if ours:
                _label(slide, column + Inches(1.70), box_top + Inches(0.42),
                       Inches(4.1), "your rover controls this one",
                       size=13, bold=True, color=TEAL, height=Inches(0.34))

    # The two sticks, tied to what they actually move.
    strip_top = row_top + (row_h + gap) * 3 + Inches(0.10)
    _box(slide, MARGIN_L, strip_top, col_w, Inches(0.58),
         "LEFT STICK on the {}  ->  surge and yaw".format(controller),
         fill=PANEL, edge=TEAL, size=15, bold=True, color=NAVY)
    _box(slide, right, strip_top, col_w, Inches(0.58),
         "RIGHT STICK  ->  the four servos",
         fill=PANEL, edge=TEAL, size=15, bold=True, color=NAVY)

    deck._note(slide,
               "A rover controls two of the six, so one stick drives it and "
               "the other is free to aim servos. A submersible controls all "
               "six, and an ROV pilot needs both sticks to fly it.",
               "info")
    return slide


# ===================================================================
# 18. plotting a maneuver from bearings
# ===================================================================

def bearing_plot(deck, title="Plotting where a maneuver ends up", speaker=None):
    """
    The trigonometry that turns a list of (time, speed, bearing) legs into a
    position on the floor.

    Redrawn from Kevin's "Running the Bases" spreadsheet. The maths is his;
    the figure is native so it prints sharp and stays editable. This is the
    slide that pays off the "trigonometry" line on the Lesson 1 STEM list.
    """
    slide = deck.blank(title, speaker=speaker or [
        "This is the trigonometry promised on the STEM slide in Lesson 1. "
        "Point back at it.",
        "Each leg of a maneuver is a distance on a bearing, and any distance "
        "on a bearing splits into an east part and a north part.",
        "Do the 45 degree case on the board first. Sine and cosine of 45 are "
        "both 0.707, so a leg at 45 degrees moves you equally in both.",
        "Bearings here are measured the navigator's way: 0 is straight ahead "
        "and they increase clockwise.",
        "Then the whole path is just adding the legs up, which is what the "
        "table on the right does.",
        "Worth saying that this is exactly how a ship's dead reckoning plot "
        "is kept - and why it drifts, because every leg's error is added to "
        "the one before it.",
    ])

    top = BODY_TOP + Inches(0.18)

    # ---------------------------------------------------------------
    # left: the formulas, and one leg drawn as a right triangle
    # ---------------------------------------------------------------
    _label(slide, MARGIN_L, top, Inches(6.0),
           "One leg: a distance on a bearing", size=16, bold=True, color=TEAL)

    _box(slide, MARGIN_L, top + Inches(0.42), Inches(6.0), Inches(0.95),
         ["distance = speed x time",
          "east  = distance x sin(bearing)",
          "north = distance x cos(bearing)"],
         fill=CODE_BG, edge=RULE, size=15, font=CODE_FONT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    # Origin is where the leg starts; the leg runs up and to the right.
    ox, oy = Inches(1.95), Inches(5.00)
    dx, dy = Inches(2.30), Inches(1.60)

    # A dashed line straight up is north, which is what a bearing is measured
    # from. Without it there is nothing for the angle to be an angle TO.
    _plain_line(slide, ox, oy, ox, oy - dy - Inches(0.30), color=GREY, width=1.0)
    _label(slide, MARGIN_L, oy - dy - Inches(0.52), Inches(1.28),
           "north", size=14, color=GREY, align=PP_ALIGN.RIGHT)
    _label(slide, MARGIN_L, oy - Inches(0.62), Inches(1.28),
           "bearing", size=14, color=GREY, align=PP_ALIGN.RIGHT)

    # The leg itself, then the two components it breaks into.
    _arrow(slide, ox, oy, ox + dx, oy - dy, color=NAVY, width=2.5)
    _plain_line(slide, ox, oy, ox + dx, oy, color=TEAL, width=1.5)
    _plain_line(slide, ox + dx, oy, ox + dx, oy - dy, color=TEAL, width=1.5)

    _label(slide, ox + Inches(0.55), oy - dy - Inches(0.22), Inches(1.7),
           "distance", size=15, bold=True, color=NAVY)
    _label(slide, ox + dx + Inches(0.14), oy - Inches(1.05), Inches(1.6),
           "north part", size=14, color=TEAL)
    _label(slide, ox + Inches(0.52), oy + Inches(0.06), Inches(1.6),
           "east part", size=14, color=TEAL)

    # ---------------------------------------------------------------
    # right: four legs, added up
    # ---------------------------------------------------------------
    tx = MARGIN_L + Inches(6.6)
    tw = Inches(5.6)
    _label(slide, tx, top, tw, "Four legs, added up as you go",
           size=16, bold=True, color=TEAL)

    headers = ["Leg", "Time", "Bearing", "East", "North"]
    rows = [
        ["1", "10 s", "45", "+23.6", "+23.6"],
        ["2", "10 s", "135", "+23.6", "-23.6"],
        ["3", "10 s", "225", "-23.6", "-23.6"],
        ["4", "10 s", "315", "-23.6", "+23.6"],
        ["", "", "total", "0.0", "0.0"],
    ]
    widths = [Inches(0.7), Inches(1.0), Inches(1.2), Inches(1.35), Inches(1.35)]

    row_h = Inches(0.46)
    head_top = top + Inches(0.46)
    x = tx
    for index, head in enumerate(headers):
        _box(slide, x, head_top, widths[index], row_h, head,
             fill=PANEL, edge=RULE, size=14, bold=True, color=NAVY,
             shape=MSO_SHAPE.RECTANGLE, edge_w=1.0)
        x += widths[index]

    for r, row in enumerate(rows):
        x = tx
        last = r == len(rows) - 1
        for index, cell in enumerate(row):
            _box(slide, x, head_top + row_h * (r + 1), widths[index], row_h,
                 cell, fill=LIGHT_TEAL if last else WHITE, edge=RULE,
                 size=14, bold=last, color=NAVY if last else INK,
                 shape=MSO_SHAPE.RECTANGLE, edge_w=1.0,
                 font=CODE_FONT if index >= 3 else BODY_FONT)
            x += widths[index]

    _label(slide, tx, head_top + row_h * (len(rows) + 1) + Inches(0.12), tw,
           "At 3.34 ft/s a 10 s leg is 33.4 ft, and sin(45) = cos(45) = 0.707.",
           size=14, color=GREY)

    deck._note(slide,
               "Four equal legs turning 90 degrees each time should end where "
               "they started. Yours will not, and the gap is the drift.",
               "info")
    return slide


# ===================================================================
# 19. the pins this vehicle uses
# ===================================================================

def pin_reference(deck, title="The ESP32 pins this vehicle actually uses",
                  motor_pins=None, led_pin="5", speaker=None):
    """
    A one-page pin reference, drawn as the board's own groups rather than as
    a full pinout.

    Kevin's Lesson 1 carries a complete ESP32-WROOM-32E pinout picture. It is
    a third-party figure and this repository is public, so this draws the
    subset that matters instead - which is also the more useful page, because
    a student looking for "which pin is the front left motor" does not want
    all 38.
    """
    slide = deck.blank(title, speaker=speaker or [
        "This is a reference page, not a teaching slide. Point at it and "
        "tell them where to find it again.",
        "The full ESP32 has 38 pins. These are the ones wired to something "
        "on this board, which is the only list they need.",
        "The motor pin numbers come straight out of the sketch, so the slide "
        "cannot disagree with the file.",
        "The two facts worth saying out loud: GPIO 34 to 39 are INPUT ONLY "
        "and have no internal pull-up, and GPIO 0 is the BOOT button, so "
        "holding it down at reset does something special.",
        "Worth printing and taping inside the lid of the parts box.",
    ])

    # srcfacts hands these back as numbers; every label on a slide is text.
    led_pin = str(led_pin)
    motor_pins = [(label, str(a), str(b)) for label, a, b in (motor_pins or [])]
    motor_pins = motor_pins or [("Front left", "12", "13"),
                                ("Rear left", "18", "19"),
                                ("Front right", "22", "23"),
                                ("Rear right", "16", "17")]

    top = BODY_TOP + Inches(0.15)
    col_w = Inches(3.95)
    gap = Inches(0.28)

    groups = [
        ("MOTORS", LIGHT_AMBER, AMBER,
         [("%s" % label, "%s / %s" % (a, b)) for label, a, b in motor_pins]
         + [("", ""), ("two pins each: drive one for", ""),
            ("forward, the other for reverse", "")]),
        ("LIGHTS AND SERVOS", LIGHT_TEAL, TEAL,
         [("32 NeoPixels", led_pin),
          ("Servo 1 / 2", "25 / 26"),
          ("Servo 3 / 4", "27 / 14"),
          ("", ""),
          ("one data wire drives all 32", ""),
          ("LEDs, in one chain", "")]),
        ("EVERYTHING ELSE", LIGHT_GREY, GREY,
         [("I2C  SDA / SCL", "32 / 33"),
          ("BOOT button", "0"),
          ("Self-test jumper", "34"),
          ("Blink LED, Lesson 1", "2"),
          ("", ""),
          ("34 to 39 are INPUT ONLY", ""),
          ("and have no pull-ups", "")]),
    ]

    for index, (heading, fill, edge, rows) in enumerate(groups):
        left = MARGIN_L + (col_w + gap) * index
        _box(slide, left, top, col_w, Inches(0.52), heading,
             fill=fill, edge=edge, size=16, bold=True, color=NAVY,
             shape=MSO_SHAPE.RECTANGLE, edge_w=1.5)

        row_top = top + Inches(0.66)
        for label, pins in rows:
            if not label:
                row_top += Inches(0.20)
                continue
            is_note = not pins
            _label(slide, left + Inches(0.06), row_top,
                   col_w - Inches(1.25) if not is_note else col_w - Inches(0.12),
                   label, size=15 if not is_note else 13,
                   color=INK if not is_note else GREY)
            if pins:
                _label(slide, left + col_w - Inches(1.15), row_top,
                       Inches(1.10), pins, size=15, bold=True, color=NAVY,
                       align=PP_ALIGN.RIGHT, font=CODE_FONT)
            row_top += Inches(0.52) if not is_note else Inches(0.34)

    deck._note(slide,
               "Every number on this page is read out of the sketches when "
               "the deck is built, so it cannot drift away from the code.",
               "info")
    return slide
