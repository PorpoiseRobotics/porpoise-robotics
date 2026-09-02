"""
slidelib.py - the slide-building helpers behind the Pathfinder course decks.

Everything here produces ordinary PowerPoint shapes: real title placeholders,
real text frames, real tables, real autoshapes. Nothing is rendered to an
image, so every word on every slide can be edited in PowerPoint afterwards.

Design constraints, in priority order:

  1. PRINTABLE. White background, dark text, no full-bleed dark fills, nothing
     that depends on colour alone to make sense in greyscale.
  2. EDITABLE. Titles live in the layout's title placeholder so the outline
     view works. Body text is in plain text frames.
  3. READABLE FROM THE BACK OF A ROOM. A hard floor on font size, and a slide
     that will not fit is SPLIT rather than shrunk into illegibility.

Slide size is 13.333 x 7.5 inches - 16:9 - which matches the existing
Porpoise Robotics decks.

FITTING
Text is measured before it is placed. If a block does not fit its box, the
font steps down towards a floor; if it still does not fit, bullets() spills
onto a continuation slide. That is why the deck builders can write as much as
the lesson needs without counting lines.
"""

import math
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ===================================================================
# PALETTE
# ===================================================================

NAVY = RGBColor(0x1F, 0x3A, 0x5F)      # Headings and rules
TEAL = RGBColor(0x0E, 0x7C, 0x86)      # Accents, "this is the point" panels
AMBER = RGBColor(0xB4, 0x6A, 0x0F)     # Warnings and watch-outs
RED = RGBColor(0xA3, 0x2A, 0x2A)       # Safety
INK = RGBColor(0x22, 0x22, 0x22)       # Body text
GREY = RGBColor(0x66, 0x66, 0x66)      # Captions, footers
RULE = RGBColor(0xC8, 0xD2, 0xDE)      # Hairlines
PANEL = RGBColor(0xEE, 0xF3, 0xF8)     # Light blue panel
PANEL_WARN = RGBColor(0xFD, 0xF3, 0xE3)
PANEL_SAFE = RGBColor(0xFB, 0xEC, 0xEC)
CODE_BG = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

BODY_FONT = "Calibri"
CODE_FONT = "Consolas"

# ===================================================================
# GEOMETRY
# ===================================================================

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN_L = Inches(0.55)
MARGIN_R = Inches(0.55)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R

TITLE_TOP = Inches(0.30)
TITLE_H = Inches(0.85)
BODY_TOP = Inches(1.30)
BODY_H = Inches(5.55)
FOOTER_TOP = Inches(6.98)

LAYOUT_TITLE = 0
LAYOUT_TITLE_CONTENT = 1
LAYOUT_SECTION = 2
LAYOUT_TITLE_ONLY = 5
LAYOUT_BLANK = 6

EMU_PER_PT = 12700

# Average advance width as a fraction of point size. Measured off Calibri and
# Consolas at 18pt; near enough for a fitting estimate, and pessimistic.
CHAR_W = {CODE_FONT: 0.56, BODY_FONT: 0.50}

# House rule: nothing a student has to READ is ever below 18pt. Code panels
# are the one exception - a listing at 18pt monospace fits barely nineteen
# lines, so they may go to 14pt to keep a whole function on one slide.
# The footer strip (track, lesson, slide number) is chrome, not content, and
# stays small.
MIN_BODY_PT = 18.0     # Body text, bullets, tables, captions, diagram labels
MIN_CODE_PT = 14.0     # Code listings only
MIN_SMALL_PT = 18.0    # Captions and diagram labels are content too
FOOTER_PT = 10.0       # Chrome


# ===================================================================
# MEASUREMENT
# ===================================================================

def _norm(runs):
    """Turns a mixed list of strings and (text, level) tuples into tuples."""
    out = []
    for item in runs:
        if isinstance(item, tuple):
            out.append((item[0], item[1]))
        else:
            out.append((item, 0))
    return out


def measure_pt(runs, width_pt, size, *, font=BODY_FONT, space_after=6,
               line_spacing=1.0):
    """Estimated rendered height of these paragraphs, in points."""
    runs = _norm(runs)
    char_w = CHAR_W.get(font, 0.50)
    total = 0.0
    for text, level in runs:
        indent = level * 18.0
        usable = max(width_pt - indent, 24.0)
        per_line = max(int(usable / (char_w * size)), 1)
        lines = max(math.ceil(len(text) / per_line), 1) if text else 1
        total += lines * size * 1.22 * line_spacing + space_after
    return total


def fit_size(runs, width_pt, height_pt, start_size, *, font=BODY_FONT,
             space_after=6, line_spacing=1.0, floor=MIN_BODY_PT):
    """
    Largest size from start_size down to floor at which these paragraphs fit
    the given box. Returns (size, space_after, fits).
    """
    size = float(start_size)
    while size >= floor:
        gap = space_after if size >= start_size - 2 else max(space_after - 3, 1)
        if measure_pt(runs, width_pt, size, font=font, space_after=gap,
                      line_spacing=line_spacing) <= height_pt:
            return size, gap, True
        size -= 1.0
    return floor, max(space_after - 4, 0), False


def _set_text(frame, runs, *, size=18, bold=False, color=INK, font=BODY_FONT,
              align=PP_ALIGN.LEFT, space_after=6, line_spacing=1.0):
    """Fills a text frame from a list of (text, level) or plain strings."""
    frame.word_wrap = True
    for index, (text, level) in enumerate(_norm(runs)):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.text = text
        para.level = min(level, 4)
        para.alignment = align
        para.space_after = Pt(space_after)
        para.line_spacing = line_spacing

        # Set the size on the PARAGRAPH as well as its runs. A blank spacer
        # line has no runs at all, and without this it takes the theme default
        # of 18pt and quietly makes the block taller than it was measured to be.
        para.font.size = Pt(size)
        para.font.name = font
        para.font.bold = bold
        para.font.color.rgb = color

        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = font
    return frame


def _set_fitted(frame, runs, *, width, height, size=18, floor=MIN_BODY_PT,
                space_after=6, **kwargs):
    """
    Fills a text frame, stepping the font down until the text fits the box.
    `width` and `height` are EMU. Returns the size actually used.
    """
    font = kwargs.get("font", BODY_FONT)
    line_spacing = kwargs.get("line_spacing", 1.0)

    inset = 2 * 7.2 if frame.margin_left is None else \
        (frame.margin_left.pt + (frame.margin_right.pt
                                 if frame.margin_right is not None else 7.2))
    width_pt = max(width / EMU_PER_PT - inset, 24.0)
    height_pt = max(height / EMU_PER_PT - 6.0, 12.0)

    used, gap, _ = fit_size(runs, width_pt, height_pt, size, font=font,
                            space_after=space_after, line_spacing=line_spacing,
                            floor=floor)
    kwargs.pop("font", None)
    kwargs.pop("line_spacing", None)
    _set_text(frame, runs, size=used, space_after=gap, font=font,
              line_spacing=line_spacing, **kwargs)
    return used


def _split_runs(runs, width_pt, height_pt, size, *, space_after=6):
    """
    Splits a list of paragraphs into chunks that each fit height_pt at `size`.
    Splits only at TOP-LEVEL items, so a sub-point never gets orphaned from
    the point it belongs to.
    """
    runs = _norm(runs)
    chunks = []
    current = []

    def height_of(block):
        return measure_pt(block, width_pt, size, space_after=space_after)

    for item in runs:
        candidate = current + [item]
        if current and item[1] == 0 and height_of(candidate) > height_pt:
            # Drop any trailing blank line before breaking.
            while current and not current[-1][0].strip():
                current.pop()
            chunks.append(current)
            current = [item]
        else:
            current = candidate

    if current:
        chunks.append(current)

    # A single top-level point with a long tail of sub-points can still be too
    # tall on its own. Break those wherever they have to break.
    final = []
    for chunk in chunks:
        if height_of(chunk) <= height_pt:
            final.append(chunk)
            continue
        part = []
        for item in chunk:
            if part and height_of(part + [item]) > height_pt:
                while part and not part[-1][0].strip():
                    part.pop()
                final.append(part)
                part = [item]
            else:
                part.append(item)
        if part:
            final.append(part)

    return _rebalance(final, runs, width_pt, height_pt, size, space_after)


def _rebalance(chunks, runs, width_pt, height_pt, size, space_after):
    """
    Uses as few slides as the content really needs, and fills them evenly.

    The greedy pass above fills each slide to the brim and spills the rest,
    which has two failure modes. It leaves a "(continued)" slide holding one
    bullet; and because it will not break inside a block of sub-points, a run
    that overflows by a line has to be hard-broken afterwards, sometimes into
    one slide more than the content needs.

    So: work out the smallest number of slides the content fits on at all,
    then aim for that many EQUAL slides. Every candidate is checked against
    the real height before it is accepted, so this can only ever produce a
    split that fits.
    """
    if len(chunks) < 2:
        return chunks

    def height_of(block):
        return measure_pt(block, width_pt, size, space_after=space_after)

    def trimmed(block):
        block = list(block)
        while block and not block[-1][0].strip():
            block.pop()
        return block

    def pack(target):
        """
        Fill to `target`, preferring to break between top-level points and
        falling back to breaking mid-block when a block is too tall on its own.
        """
        out, current = [], []
        for item in runs:
            candidate = current + [item]
            if current and height_of(candidate) > target:
                if item[1] == 0 or height_of(current) > height_pt:
                    out.append(trimmed(current))
                    current = [item]
                    continue
            current = candidate
        if current:
            out.append(trimmed(current))
        return [block for block in out if block]

    def usable(packed, want):
        return (len(packed) == want
                and all(height_of(block) <= height_pt for block in packed))

    total = height_of(runs)
    best = None
    for count in range(2, len(chunks) + 1):
        # Walk the target from "perfectly even" up to a full slide, and keep
        # the evenest packing that still lands on exactly `count` slides.
        steps = 14
        for step in range(steps + 1):
            target = total / count
            target += (height_pt - target) * step / steps
            if target > height_pt:
                break
            packed = pack(target)
            if not usable(packed, count):
                continue
            heights = [height_of(block) for block in packed]
            spread = max(heights) - min(heights)
            if best is None or spread < best[0]:
                best = (spread, packed)
        if best is not None:
            return best[1]

    return chunks


class Placeholder:
    """
    A picture we have not taken yet.

    Anywhere a slide kind takes an image path it will also take one of these,
    and it will also take the path of a photograph that does not exist. Either
    way you get a labelled dashed box on the slide, in the space the picture
    will occupy, saying what belongs there - so a gap in the artwork is
    visible on the slide rather than silently missing from it.

        deck.image_slide("The circuit", Placeholder(
            "PHOTO: our breadboard, LED and 220 ohm resistor, wired to GPIO 2",
            "Kevin's photographs of this use a 9 V battery instead."))
    """

    def __init__(self, label, hint=None):
        self.label = label
        self.hint = hint


def set_notes(slide, notes):
    """
    Writes the teacher's script into the slide's notes page.

    `notes` may be a string or a list of lines. These are what a teacher sees
    in Presenter View and on a printed notes page; students never see them.
    """
    if not notes:
        return
    if isinstance(notes, str):
        notes = [notes]
    frame = slide.notes_slide.notes_text_frame
    frame.text = ""
    for index, line in enumerate(notes):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.text = line
        for run in para.runs:
            run.font.size = Pt(12)
            run.font.name = BODY_FONT
    return slide


def set_notes_all(slides, notes):
    """
    The same script on every slide of a group.

    bullets(), code(), two_columns() and activity() all spill onto
    "(continued)" slides when the content is too tall. A teacher who advances
    onto one of those should not find the notes pane empty, so each carries
    the same script with a line at the top saying it is a continuation.
    """
    if not notes:
        return slides[0] if slides else None
    if isinstance(notes, str):
        notes = [notes]
    for index, slide in enumerate(slides):
        if index == 0:
            set_notes(slide, notes)
        else:
            set_notes(slide, ["(Continued from the previous slide - the same "
                              "notes apply.)", ""] + list(notes))
    return slides[0]


class Deck:
    """One lesson deck."""

    def __init__(self, filename, deck_title, lesson_label, track_label,
                 footer_note="Porpoise Robotics - proprietary. For use with "
                             "Porpoise Robotics lessons."):
        self.filename = filename
        self.deck_title = deck_title
        self.lesson_label = lesson_label
        self.track_label = track_label
        self.footer_note = footer_note

        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.slide_number = 0

    # ---------------------------------------------------------------
    # plumbing
    # ---------------------------------------------------------------

    def _new(self, layout=LAYOUT_TITLE_ONLY, title=None, numbered=True):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout])

        for shape in list(slide.placeholders):
            if shape.placeholder_format.idx != 0:
                shape._element.getparent().remove(shape._element)

        if title is not None and slide.shapes.title is not None:
            holder = slide.shapes.title
            holder.left = MARGIN_L
            holder.top = TITLE_TOP
            holder.width = CONTENT_W
            holder.height = TITLE_H
            frame = holder.text_frame
            frame.word_wrap = True
            frame.vertical_anchor = MSO_ANCHOR.BOTTOM
            _set_fitted(frame, [title], width=CONTENT_W, height=TITLE_H,
                        size=28, floor=17, bold=True, color=NAVY, space_after=0)
            self._rule(slide, TITLE_TOP + TITLE_H + Inches(0.05))

        if numbered:
            self.slide_number += 1
            self._footer(slide)
        return slide

    def _rule(self, slide, top):
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, top,
                                      CONTENT_W, Pt(1.5))
        line.fill.solid()
        line.fill.fore_color.rgb = RULE
        line.line.fill.background()
        line.shadow.inherit = False
        return line

    def _footer(self, slide):
        box = slide.shapes.add_textbox(MARGIN_L, FOOTER_TOP,
                                       CONTENT_W - Inches(0.7), Inches(0.3))
        frame = box.text_frame
        frame.margin_left = 0
        frame.margin_top = 0
        _set_text(frame, ["{}  |  {}".format(self.track_label, self.lesson_label)],
                  size=FOOTER_PT, color=GREY, space_after=0)

        number = slide.shapes.add_textbox(SLIDE_W - Inches(1.05), FOOTER_TOP,
                                          Inches(0.5), Inches(0.3))
        nframe = number.text_frame
        nframe.margin_left = 0
        nframe.margin_top = 0
        _set_text(nframe, [str(self.slide_number)], size=FOOTER_PT, color=GREY,
                  align=PP_ALIGN.RIGHT, space_after=0)

    def _panel(self, slide, left, top, width, height, fill, edge=None):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                     width, height)
        box.adjustments[0] = 0.04
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        if edge is None:
            box.line.fill.background()
        else:
            box.line.color.rgb = edge
            box.line.width = Pt(1)
        box.shadow.inherit = False
        box.text_frame.word_wrap = True
        # python-pptx builds an autoshape with algn="ctr" on its first
        # paragraph. Left is what we want everywhere unless a caller says
        # otherwise, and a centred line of code looks broken.
        box.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        return box

    def _place_image(self, slide, image, left, top, width, height):
        """
        A picture if we have one, a labelled dashed box if we do not.

        Returns the shape that was placed, or None if `image` was None. A
        picture keeps its aspect ratio inside the (width, height) box given
        and is centred in it; a placeholder fills the box exactly, so the
        caption below it lands where the caption of the real photograph will.
        """
        if image is None:
            return None

        if isinstance(image, Placeholder):
            return self._placeholder_box(slide, left, top, width, height,
                                         image.label, image.hint)

        if not os.path.exists(image):
            # A path that does not resolve is a picture somebody meant to
            # supply. Say so on the slide rather than leaving a hole.
            return self._placeholder_box(
                slide, left, top, width, height,
                "IMAGE TO COME",
                "Expected file: " + os.path.basename(str(image)))

        pic = slide.shapes.add_picture(image, left, top, width=width)
        if pic.height > height:
            scale = height / pic.height
            pic.height = int(pic.height * scale)
            pic.width = int(pic.width * scale)
        pic.left = left + Emu(int((width - pic.width) / 2))
        pic.top = top + Emu(int((height - pic.height) / 2))
        return pic

    def _placeholder_box(self, slide, left, top, width, height, label,
                         hint=None):
        """
        The dashed box that stands in for a picture we have not got yet.

        Deliberately plain and deliberately obvious: a teacher flipping
        through the deck should be able to spot every outstanding photograph
        in one pass, and a printed copy should show them too.
        """
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top,
                                     width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = PANEL
        box.line.color.rgb = TEAL
        box.line.width = Pt(1.5)
        box.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        box.shadow.inherit = False

        rows = [("PLACEHOLDER", 0), (label, 0)]
        if hint:
            rows.append((hint, 0))

        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = Inches(0.18)
        frame.margin_right = Inches(0.18)
        frame.margin_top = Inches(0.12)
        frame.margin_bottom = Inches(0.12)
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_fitted(frame, rows, width=width, height=height, size=16,
                    floor=11, color=TEAL, align=PP_ALIGN.CENTER,
                    space_after=7)
        # "PLACEHOLDER" reads as a label, so set it apart from the sentence
        # underneath it that says what the picture will be.
        head = frame.paragraphs[0]
        for run in head.runs:
            run.font.bold = True
            run.font.color.rgb = GREY
        return box

    def _textbox(self, slide, left, top, width, height):
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.word_wrap = True
        frame.margin_left = 0
        frame.margin_top = 0
        frame.margin_right = 0
        frame.margin_bottom = 0
        return box

    # ---------------------------------------------------------------
    # slide kinds
    # ---------------------------------------------------------------

    def title_slide(self, subtitle, byline_lines, hero_image=None, logo=None,
                    speaker=None):
        slide = self._new(layout=LAYOUT_BLANK, title=None, numbered=False)

        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                      SLIDE_W, Inches(0.22))
        band.fill.solid()
        band.fill.fore_color.rgb = NAVY
        band.line.fill.background()
        band.shadow.inherit = False

        # House mark, top right, clear of the hero image below it.
        if logo and os.path.exists(logo):
            mark = slide.shapes.add_picture(logo, 0, Inches(0.5),
                                            height=Inches(0.85))
            mark.left = SLIDE_W - MARGIN_R - mark.width

        text_w = Inches(7.1) if hero_image else CONTENT_W

        box = self._textbox(slide, MARGIN_L, Inches(1.15), text_w, Inches(0.4))
        _set_text(box.text_frame, [self.lesson_label.upper()], size=16, bold=True,
                  color=TEAL, space_after=0)

        box = self._textbox(slide, MARGIN_L, Inches(1.62), text_w, Inches(1.7))
        _set_fitted(box.text_frame, [self.deck_title], width=text_w,
                    height=Inches(1.7), size=40, floor=24, bold=True,
                    color=NAVY, space_after=0, line_spacing=0.95)

        box = self._textbox(slide, MARGIN_L, Inches(3.45), text_w, Inches(1.05))
        _set_fitted(box.text_frame, [subtitle], width=text_w,
                    height=Inches(1.05), size=17, floor=13, color=INK,
                    space_after=0)

        # The credit block is chrome, like the footer - it names the team
        # rather than teaching anything, so it is exempt from the 18pt floor.
        box = self._textbox(slide, MARGIN_L, Inches(4.30), text_w, Inches(2.35))
        _set_fitted(box.text_frame, byline_lines, width=text_w,
                    height=Inches(2.35), size=13, floor=10, color=GREY,
                    space_after=1)

        if hero_image is not None:
            self._place_image(slide, hero_image, Inches(7.9), Inches(1.5),
                              Inches(4.9), Inches(5.1))

        box = self._textbox(slide, MARGIN_L, Inches(6.72), CONTENT_W, Inches(0.55))
        _set_text(box.text_frame, [self.footer_note], size=9, color=GREY,
                  space_after=0)
        set_notes(slide, speaker)
        return slide

    def section(self, title, subtitle=None, minutes=None, speaker=None):
        slide = self._new(layout=LAYOUT_BLANK, title=None)

        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.4),
                                     Inches(0.28), Inches(2.0))
        bar.fill.solid()
        bar.fill.fore_color.rgb = TEAL
        bar.line.fill.background()
        bar.shadow.inherit = False

        box = self._textbox(slide, Inches(0.85), Inches(2.5), Inches(11.5), Inches(1.3))
        _set_fitted(box.text_frame, [title], width=Inches(11.5),
                    height=Inches(1.3), size=36, floor=22, bold=True,
                    color=NAVY, space_after=0)

        rows = []
        if subtitle:
            rows.append(subtitle)
        if minutes:
            rows.append("About {} minutes".format(minutes))
        if rows:
            box = self._textbox(slide, Inches(0.85), Inches(3.85), Inches(11.5),
                                Inches(1.0))
            _set_fitted(box.text_frame, rows, width=Inches(11.5),
                        height=Inches(1.0), size=16, floor=12, color=GREY,
                        space_after=4)
        set_notes(slide, speaker)
        return slide

    def objectives(self, items, title="What you will be able to do by the end",
                   speaker=None):
        """
        The learning objectives for the lesson, stated as things the student
        will be able to DO. Shown near the front and worth returning to at the
        end.
        """
        slide = self._new(title=title)

        box = self._textbox(slide, MARGIN_L, BODY_TOP, CONTENT_W, Inches(0.55))
        _set_fitted(box.text_frame, ["By the end of this lesson you will be "
                                     "able to:"], width=CONTENT_W,
                    height=Inches(0.55), size=19, bold=True, color=TEAL)

        numbered = [("%d.   %s" % (i + 1, text), 0)
                    for i, text in enumerate(items)]
        height = BODY_H - Inches(0.6)
        box = self._textbox(slide, MARGIN_L, BODY_TOP + Inches(0.6), CONTENT_W,
                            height)
        _set_fitted(box.text_frame, numbered, width=CONTENT_W, height=height,
                    size=19, space_after=13)

        set_notes(slide, speaker)
        return slide

    def progress(self, stages, done, title="Where we are", speaker=None):
        """
        A strip of the lesson's stages with the finished ones ticked and the
        current one highlighted. Shown between sections so the class can see
        how far through the three hours they are.
        """
        slide = self._new(title=title)

        gap = Inches(0.16)
        count = len(stages)
        box_w = Emu(int((CONTENT_W - gap * (count - 1)) / count))
        top = BODY_TOP + Inches(1.15)
        box_h = Inches(1.9)

        for index, stage in enumerate(stages):
            left = MARGIN_L + (box_w + gap) * index
            if index < done:
                fill, edge, mark = PANEL, TEAL, "done"
            elif index == done:
                fill, edge, mark = RGBColor(0xD6, 0xEC, 0xEE), TEAL, "NOW"
            else:
                fill, edge, mark = WHITE, RULE, ""

            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left,
                                           top, box_w, box_h)
            shape.adjustments[0] = 0.08
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
            shape.line.color.rgb = edge
            shape.line.width = Pt(2.0 if index == done else 1.0)
            shape.shadow.inherit = False

            frame = shape.text_frame
            frame.word_wrap = True
            frame.margin_left = Inches(0.08)
            frame.margin_right = Inches(0.08)
            frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            _set_fitted(frame, [stage], width=box_w, height=box_h, size=17,
                        floor=13, bold=(index == done),
                        color=NAVY if index <= done else GREY,
                        align=PP_ALIGN.CENTER, space_after=0)

            if mark:
                tag = self._textbox(slide, left, top - Inches(0.42), box_w,
                                    Inches(0.38))
                _set_fitted(tag.text_frame, [mark], width=box_w,
                            height=Inches(0.38), size=13, floor=11, bold=True,
                            color=TEAL, align=PP_ALIGN.CENTER)

        set_notes(slide, speaker)
        return slide

    def bullets(self, title, items, lead=None, note=None, note_kind="info",
                size=18, speaker=None):
        """
        A bulleted slide. If the content will not fit even at the minimum font
        size, it spills onto "(continued)" slides, splitting only between
        top-level points.
        """
        top = BODY_TOP
        height = BODY_H
        if lead:
            height -= Inches(0.62)
        if note:
            height -= Inches(1.15)

        width_pt = CONTENT_W / EMU_PER_PT
        # _set_fitted keeps 6pt of slack inside the box. Measure against the
        # same figure, or we decide "it fits" and then overflow by a hair.
        height_pt = height / EMU_PER_PT - 8.0

        # The gap between paragraphs has to match what _set_fitted will use,
        # or the split decision is made against the wrong height.
        gap = 9

        # Would it fit if we shrank it? If so, one slide will do.
        _, _, fits = fit_size(items, width_pt, height_pt, size, space_after=gap)
        if fits:
            chunks = [_norm(items)]
        else:
            chunks = _split_runs(items, width_pt, height_pt, MIN_BODY_PT,
                                 space_after=gap)

        slides = []
        for index, chunk in enumerate(chunks):
            slide_title = title if index == 0 else title + "  (continued)"
            slide = self._new(title=slide_title)
            slides.append(slide)

            body_top = BODY_TOP
            body_h = BODY_H
            is_last = index == len(chunks) - 1

            if lead and index == 0:
                box = self._textbox(slide, MARGIN_L, body_top, CONTENT_W,
                                    Inches(0.58))
                _set_fitted(box.text_frame, [lead], width=CONTENT_W,
                            height=Inches(0.58), size=19, floor=14, bold=True,
                            color=TEAL)
                body_top += Inches(0.62)
                body_h -= Inches(0.62)

            if note and is_last:
                body_h -= Inches(1.15)

            box = self._textbox(slide, MARGIN_L, body_top, CONTENT_W, body_h)
            _set_fitted(box.text_frame, chunk, width=CONTENT_W, height=body_h,
                        size=size, space_after=9)

            if note and is_last:
                self._note(slide, note, note_kind)

        set_notes_all(slides, speaker)
        return slides[0]

    def _note(self, slide, text, kind="info", top=None, left=None, width=None):
        fills = {"info": (PANEL, TEAL), "warn": (PANEL_WARN, AMBER),
                 "safety": (PANEL_SAFE, RED)}
        labels = {"info": "KEY POINT", "warn": "WATCH OUT", "safety": "SAFETY"}
        fill, edge = fills.get(kind, fills["info"])
        label = labels.get(kind, "KEY POINT")

        left = MARGIN_L if left is None else left
        width = CONTENT_W if width is None else width

        # Grow the panel if the note is long, rather than overflowing it.
        width_pt = width / EMU_PER_PT - 0.5
        combined = label + "   " + text
        size, _, _ = fit_size([combined], width_pt, 1.0 * 72 - 12, 15, floor=12)
        wanted_pt = measure_pt([combined], width_pt, size, space_after=0)
        height = max(Inches(1.05), Emu(int((wanted_pt + 18) * EMU_PER_PT)))
        height = min(height, Inches(1.7))

        top = Inches(6.85) - height if top is None else top

        panel = self._panel(slide, left, top, width, height, fill, edge)
        frame = panel.text_frame
        frame.margin_left = Inches(0.18)
        frame.margin_right = Inches(0.18)
        frame.margin_top = Inches(0.08)
        frame.margin_bottom = Inches(0.08)
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        para = frame.paragraphs[0]
        run = para.add_run()
        run.text = label + "   "
        run.font.size = Pt(min(size, 13))
        run.font.bold = True
        run.font.color.rgb = edge
        run.font.name = BODY_FONT

        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = INK
        run.font.name = BODY_FONT
        return panel

    def bullets_image(self, title, items, image, caption=None, note=None,
                      note_kind="info", image_ratio=0.42, size=17,
                      speaker=None):
        slide = self._new(title=title)

        img_w = Emu(int(CONTENT_W * image_ratio))
        txt_w = CONTENT_W - img_w - Inches(0.4)

        body_h = BODY_H - (Inches(1.25) if note else Inches(0))
        box = self._textbox(slide, MARGIN_L, BODY_TOP, txt_w, body_h)
        _set_fitted(box.text_frame, items, width=txt_w, height=body_h,
                    size=size, space_after=9)

        # A photograph is usually shorter than the box it is given, so its
        # caption floats up with it. A placeholder fills the box exactly, so
        # the space for the caption - and for the note panel below it - has
        # to be taken out of the box first rather than borrowed afterwards.
        cap_h = Inches(0.62) if caption else Inches(0)
        img_bottom = Inches(5.52) if note else Inches(6.85)
        img_h = min(Inches(4.2), img_bottom - BODY_TOP - cap_h)

        pic = self._place_image(slide, image, MARGIN_L + txt_w + Inches(0.4),
                                BODY_TOP, img_w, img_h)
        if pic is not None and caption:
            cap = self._textbox(slide, pic.left,
                                pic.top + pic.height + Inches(0.06),
                                pic.width, cap_h)
            _set_fitted(cap.text_frame, [caption], width=pic.width,
                        height=cap_h, size=MIN_SMALL_PT,
                        floor=MIN_SMALL_PT, color=GREY)

        if note:
            self._note(slide, note, note_kind)
        set_notes(slide, speaker)
        return slide

    def image_slide(self, title, image, caption=None, items=None, note=None,
                    note_kind="info", speaker=None):
        slide = self._new(title=title)

        top = BODY_TOP
        if items:
            box = self._textbox(slide, MARGIN_L, top, CONTENT_W, Inches(1.1))
            _set_fitted(box.text_frame, items, width=CONTENT_W,
                        height=Inches(1.1), size=16, space_after=6)
            top += Inches(1.18)

        bottom_limit = Inches(5.55) if note else Inches(6.85)
        available_h = bottom_limit - top - (Inches(0.80) if caption else Inches(0))

        if available_h > Inches(0.6):
            pic = self._place_image(slide, image, MARGIN_L, top,
                                    CONTENT_W, available_h)

            if pic is not None and caption:
                cap = self._textbox(slide, MARGIN_L,
                                    pic.top + pic.height + Inches(0.06),
                                    CONTENT_W, Inches(0.74))
                _set_fitted(cap.text_frame, [caption], width=CONTENT_W,
                            height=Inches(0.74), size=MIN_SMALL_PT,
                            floor=MIN_SMALL_PT, color=GREY, align=PP_ALIGN.CENTER)

        if note:
            self._note(slide, note, note_kind)
        set_notes(slide, speaker)
        return slide

    def image_pair(self, title, left_image, left_caption, right_image,
                   right_caption, lead=None, note=None, note_kind="info",
                   speaker=None):
        """Two photographs side by side, for a before/after or left/right pair."""
        slide = self._new(title=title)

        top = BODY_TOP
        if lead:
            box = self._textbox(slide, MARGIN_L, top, CONTENT_W, Inches(0.6))
            _set_fitted(box.text_frame, [lead], width=CONTENT_W,
                        height=Inches(0.6), size=MIN_BODY_PT, color=INK)
            top += Inches(0.68)

        gap = Inches(0.4)
        col_w = Emu(int((CONTENT_W - gap) / 2))
        bottom = Inches(5.55) if note else Inches(6.85)
        pic_h = bottom - top - Inches(0.85)

        pairs = ((left_image, left_caption), (right_image, right_caption))

        # Photographs go down first. A photograph keeps its aspect ratio and
        # is usually shorter than the box; a placeholder fills whatever box it
        # is given. Sizing the placeholders to the tallest real picture keeps
        # a mixed pair looking like a pair rather than a picture next to a
        # much bigger empty rectangle.
        placed = {}
        for index, (path, _) in enumerate(pairs):
            if path is None or isinstance(path, Placeholder):
                continue
            left = MARGIN_L + (col_w + gap) * index
            placed[index] = self._place_image(slide, path, left, top,
                                              col_w, pic_h)

        real = [pic for pic in placed.values() if pic is not None]
        box_h = max((pic.height for pic in real), default=pic_h)

        for index, (path, caption) in enumerate(pairs):
            if index not in placed:
                if path is None:
                    continue
                left = MARGIN_L + (col_w + gap) * index
                placed[index] = self._place_image(slide, path, left, top,
                                                  col_w, box_h)

            pic = placed[index]
            if pic is None:
                continue

            left = MARGIN_L + (col_w + gap) * index
            cap = self._textbox(slide, left, top + box_h + Inches(0.08),
                                col_w, Inches(0.72))
            _set_fitted(cap.text_frame, [caption], width=col_w,
                        height=Inches(0.72), size=MIN_SMALL_PT,
                        floor=MIN_SMALL_PT, color=GREY, align=PP_ALIGN.CENTER)

        if note:
            self._note(slide, note, note_kind)
        set_notes(slide, speaker)
        return slide

    def two_columns(self, title, left_heading, left_items, right_heading,
                    right_items, note=None, note_kind="info", size=16,
                    speaker=None):
        """
        Two headed columns. If either column will not fit at the minimum font
        size, BOTH spill onto a continuation slide together, so the pairing
        stays readable.
        """
        gap = Inches(0.5)
        col_w = Emu(int((CONTENT_W - gap) / 2))
        body_h = BODY_H - (Inches(1.25) if note else Inches(0)) - Inches(0.5)

        width_pt = col_w / EMU_PER_PT
        height_pt = body_h / EMU_PER_PT - 8.0

        _, _, left_fits = fit_size(left_items, width_pt, height_pt, size,
                                   space_after=8)
        _, _, right_fits = fit_size(right_items, width_pt, height_pt, size,
                                    space_after=8)

        if left_fits and right_fits:
            left_chunks = [_norm(left_items)]
            right_chunks = [_norm(right_items)]
        else:
            left_chunks = _split_runs(left_items, width_pt, height_pt,
                                      MIN_BODY_PT, space_after=8)
            right_chunks = _split_runs(right_items, width_pt, height_pt,
                                       MIN_BODY_PT, space_after=8)

        pages = max(len(left_chunks), len(right_chunks))
        left_chunks += [[]] * (pages - len(left_chunks))
        right_chunks += [[]] * (pages - len(right_chunks))

        first = None
        made = []
        for index in range(pages):
            slide_title = title if index == 0 else title + "  (continued)"
            slide = self._new(title=slide_title)
            first = first or slide
            made.append(slide)
            is_last = index == pages - 1

            this_h = BODY_H - (Inches(1.25) if (note and is_last) else Inches(0))

            for column, (heading, chunk) in enumerate(
                    ((left_heading, left_chunks[index]),
                     (right_heading, right_chunks[index]))):
                left = MARGIN_L + (col_w + gap) * column
                if not chunk and index > 0:
                    continue

                head = self._textbox(slide, left, BODY_TOP, col_w, Inches(0.5))
                _set_fitted(head.text_frame,
                            [heading if index == 0 else heading + " (cont.)"],
                            width=col_w, height=Inches(0.5), size=18,
                            floor=MIN_BODY_PT, bold=True, color=TEAL)

                box = self._textbox(slide, left, BODY_TOP + Inches(0.56), col_w,
                                    this_h - Inches(0.56))
                _set_fitted(box.text_frame, chunk, width=col_w,
                            height=this_h - Inches(0.56), size=size,
                            space_after=8)

            if note and is_last:
                self._note(slide, note, note_kind)

        set_notes_all(made, speaker)
        return first

    def code(self, title, lines, filename=None, notes=None, size=13,
             highlight=None, speaker=None):
        """
        A code slide. `lines` is the listing. `notes` is an optional list of
        short explanations shown down the right-hand side. `highlight` is a
        set of zero-based line numbers to draw in the accent colour.

        A listing too long for one panel even at the code floor is split over
        continuation slides; the explanations stay with the first.
        """
        highlight = set(highlight or ())

        code_w = CONTENT_W if not notes else Emu(int(CONTENT_W * 0.63))
        head_h = Inches(0.36) if filename else Inches(0)
        panel_h = Inches(6.85) - (BODY_TOP + head_h)

        # How many lines fit at the code floor, and how wide they may be.
        avail_w_pt = code_w / EMU_PER_PT - 20
        avail_h_pt = panel_h / EMU_PER_PT - 14
        longest = max((len(l) for l in lines), default=1)

        code_size = float(size)
        while code_size > MIN_CODE_PT:
            if (longest * CHAR_W[CODE_FONT] * code_size <= avail_w_pt and
                    len(lines) * code_size * 1.22 <= avail_h_pt):
                break
            code_size -= 0.5

        per_page = max(int(avail_h_pt / (code_size * 1.22)), 4)
        pages = [lines[i:i + per_page] for i in range(0, len(lines), per_page)]             or [lines]

        first = None
        made = []
        for index, chunk in enumerate(pages):
            slide = self._new(title=title if index == 0
                              else title + "  (continued)")
            first = first or slide
            made.append(slide)
            top = BODY_TOP

            if filename:
                box = self._textbox(slide, MARGIN_L, top, code_w, Inches(0.3))
                _set_fitted(box.text_frame, [filename], width=code_w,
                            height=Inches(0.3), size=14, floor=12,
                            bold=True, color=TEAL, font=CODE_FONT)
                top += Inches(0.36)

            this_h = Inches(6.85) - top
            panel = self._panel(slide, MARGIN_L, top, code_w, this_h,
                                CODE_BG, RULE)
            panel.name = "CODE_PANEL"
            frame = panel.text_frame
            frame.margin_left = Inches(0.16)
            frame.margin_right = Inches(0.1)
            frame.margin_top = Inches(0.1)
            frame.margin_bottom = Inches(0.08)
            frame.vertical_anchor = MSO_ANCHOR.TOP

            offset = index * per_page
            for row, line in enumerate(chunk):
                para = frame.paragraphs[0] if row == 0 else frame.add_paragraph()
                para.text = line if line else " "
                para.space_after = Pt(0)
                para.line_spacing = 1.0
                para.alignment = PP_ALIGN.LEFT
                para.font.size = Pt(code_size)
                para.font.name = CODE_FONT
                stripped = line.strip()
                for run in para.runs:
                    run.font.size = Pt(code_size)
                    run.font.name = CODE_FONT
                    if (offset + row) in highlight:
                        run.font.bold = True
                        run.font.color.rgb = TEAL
                    elif stripped.startswith(("//", "/*", "*")):
                        run.font.color.rgb = GREY
                    else:
                        run.font.color.rgb = INK

            # Explanations belong beside the first page of the listing.
            if notes and index == 0:
                note_left = MARGIN_L + code_w + Inches(0.35)
                note_w = CONTENT_W - code_w - Inches(0.35)
                box = self._textbox(slide, note_left, top, note_w, this_h)
                _set_fitted(box.text_frame, notes, width=note_w, height=this_h,
                            size=16, space_after=10)

        set_notes_all(made, speaker)
        return first

    def table(self, title, headers, rows, lead=None, note=None,
              note_kind="info", col_widths=None, size=14, speaker=None):
        slide = self._new(title=title)
        top = BODY_TOP

        if lead:
            box = self._textbox(slide, MARGIN_L, top, CONTENT_W, Inches(0.5))
            _set_fitted(box.text_frame, [lead], width=CONTENT_W,
                        height=Inches(0.5), size=16, floor=12, color=INK)
            top += Inches(0.55)

        bottom = Inches(5.55) if note else Inches(6.85)
        available = bottom - top

        # Estimate how tall each row wants to be, so the table does not spill.
        widths = col_widths or [1] * len(headers)
        total_weight = float(sum(widths))
        col_pt = [(CONTENT_W / EMU_PER_PT) * w / total_weight - 10
                  for w in widths]

        cell_size = float(size)
        while cell_size >= 9.0:
            needed = 0.0
            for row in [headers] + [list(map(str, r)) for r in rows]:
                tallest = 0.0
                for text, width_pt in zip(row, col_pt):
                    per_line = max(int(width_pt / (CHAR_W[BODY_FONT] * cell_size)), 1)
                    lines = max(math.ceil(len(text) / per_line), 1)
                    tallest = max(tallest, lines * cell_size * 1.25 + 6)
                needed += tallest
            if needed <= available / EMU_PER_PT:
                break
            cell_size -= 0.5

        height = min(available, Emu(int(needed * EMU_PER_PT)))
        shape = slide.shapes.add_table(len(rows) + 1, len(headers), MARGIN_L,
                                       top, CONTENT_W, height)
        table = shape.table
        table.first_row = True

        if col_widths:
            for index, weight in enumerate(col_widths):
                table.columns[index].width = Emu(int(CONTENT_W * weight / total_weight))

        for col, text in enumerate(headers):
            cell = table.cell(0, col)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            cell.margin_left = Inches(0.08)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(cell_size)
                    run.font.bold = True
                    run.font.color.rgb = WHITE
                    run.font.name = BODY_FONT

        for r, row in enumerate(rows, start=1):
            for c, text in enumerate(row):
                cell = table.cell(r, c)
                cell.text = str(text)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 else PANEL
                cell.margin_left = Inches(0.08)
                cell.margin_top = Inches(0.02)
                cell.margin_bottom = Inches(0.02)
                mono = c == 0 and any(ch in str(text) for ch in "_(")
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(cell_size)
                        run.font.color.rgb = INK
                        run.font.name = CODE_FONT if mono else BODY_FONT

        if note:
            self._note(slide, note, note_kind)
        set_notes(slide, speaker)
        return slide

    def activity(self, title, sketch, steps, expect=None, questions=None,
                 minutes=None, safety=None, speaker=None):
        """
        A 'stop and do this' slide, deliberately different from the rest.
        Long step lists spill onto a continuation slide; what students should
        see and the questions stay with the first one.
        """
        gap = Inches(0.45)
        has_right = bool(expect or questions)
        col_w = Emu(int((CONTENT_W - gap) / 2)) if has_right else CONTENT_W

        top = BODY_TOP + Inches(0.78)
        bottom = Inches(5.55) if safety else Inches(6.85)
        col_h = bottom - top

        width_pt = col_w / EMU_PER_PT
        height_pt = col_h / EMU_PER_PT - 8.0

        _, _, fits = fit_size(steps, width_pt, height_pt, 16, space_after=8)
        chunks = [_norm(steps)] if fits else _split_runs(
            steps, width_pt, height_pt, MIN_BODY_PT, space_after=8)

        first = None
        made = []
        for index, chunk in enumerate(chunks):
            slide = self._new(title=title if index == 0
                              else title + "  (continued)")
            first = first or slide
            made.append(slide)
            is_last = index == len(chunks) - 1

            header = self._panel(slide, MARGIN_L, BODY_TOP, CONTENT_W,
                                 Inches(0.58), TEAL)
            frame = header.text_frame
            frame.margin_left = Inches(0.2)
            frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            para = frame.paragraphs[0]

            run = para.add_run()
            run.text = "UPLOAD AND RUN:   "
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = WHITE
            run.font.name = BODY_FONT

            run = para.add_run()
            run.text = sketch
            run.font.size = Pt(15)
            run.font.bold = True
            run.font.color.rgb = WHITE
            run.font.name = CODE_FONT

            if minutes and index == 0:
                run = para.add_run()
                run.text = "     ({} minutes)".format(minutes)
                run.font.size = Pt(12)
                run.font.color.rgb = WHITE
                run.font.name = BODY_FONT

            this_bottom = Inches(5.55) if (safety and is_last) else Inches(6.85)
            this_h = this_bottom - top

            box = self._textbox(slide, MARGIN_L, top, col_w, this_h)
            _set_fitted(box.text_frame, chunk, width=col_w, height=this_h,
                        size=16, space_after=8)

            # What to look for and the questions belong with the first page.
            if has_right and index == 0:
                right = MARGIN_L + col_w + gap
                right_top = top
                blocks = []
                if expect:
                    blocks.append(("What you should see", expect))
                if questions:
                    blocks.append(("Answer these", questions))

                share = Emu(int((this_h - Inches(0.44) * len(blocks))
                                / len(blocks)))
                for heading, body in blocks:
                    head = self._textbox(slide, right, right_top, col_w,
                                         Inches(0.4))
                    _set_fitted(head.text_frame, [heading], width=col_w,
                                height=Inches(0.4), size=16,
                                floor=MIN_BODY_PT, bold=True, color=TEAL)
                    box = self._textbox(slide, right, right_top + Inches(0.44),
                                        col_w, share)
                    _set_fitted(box.text_frame, body, width=col_w, height=share,
                                size=15, space_after=6)
                    right_top += share + Inches(0.48)

            if safety and is_last:
                self._note(slide, safety, "safety")

        set_notes_all(made, speaker)
        return first

    def quiz(self, title, questions, lead=None, speaker=None):
        slide = self._new(title=title)
        top = BODY_TOP
        if lead:
            box = self._textbox(slide, MARGIN_L, top, CONTENT_W, Inches(0.45))
            _set_fitted(box.text_frame, [lead], width=CONTENT_W,
                        height=Inches(0.45), size=16, floor=12, color=GREY)
            top += Inches(0.52)

        height = Inches(6.85) - top
        box = self._textbox(slide, MARGIN_L, top, CONTENT_W, height)
        _set_fitted(box.text_frame, questions, width=CONTENT_W, height=height,
                    size=17, space_after=14)
        set_notes(slide, speaker)
        return slide

    def blank(self, title, speaker=None):
        slide = self._new(title=title)
        set_notes(slide, speaker)
        return slide

    # ---------------------------------------------------------------
    # output
    # ---------------------------------------------------------------

    def save(self, folder):
        path = os.path.join(folder, self.filename)
        self.prs.save(path)
        return path, len(self.prs.slides._sldIdLst)
